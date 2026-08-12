"""Builds Analysis_Module_Overview.pptx -- a teaching deck for Module 2 (Analysis).

Every number in this deck comes from actually running analyze.py against
Module 1's real output (3 documents, as of this writing) -- not illustrative
placeholders. Re-run analyze.py first if you want the deck to reflect a
fresh Module 1 run; this script doesn't re-run it for you.

Same visual design system as 1. OCR/make_slides.py, duplicated here rather
than imported -- each module stays self-contained and only shares the venv,
same convention as tokenizer_util.py existing separately in both modules.

Usage:
  python make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "Analysis_Module_Overview.pptx"

# ---- palette (matches 1. OCR/make_slides.py) --------------------------------
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x27)
MUTED = RGBColor(0x8A, 0x93, 0xA6)
MUTED_DARK = RGBColor(0x5B, 0x63, 0x74)
ACCENT = RGBColor(0x3E, 0x8E, 0xDE)       # blue -- structural content
ACCENT_BAD = RGBColor(0xE0, 0x6C, 0x5C)   # red -- real findings / root cause analysis

FONT_BODY = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def _kicker(slide, text, color=ACCENT):
    _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5), text.upper(), 15, color, bold=True)


def _accent_bar(slide, color, top=Inches(0), height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _page_number(slide, n, dark=False):
    color = MUTED if dark else MUTED_DARK
    _textbox(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35), str(n), 11, color, align=PP_ALIGN.RIGHT)


# ---- adaptive bullet sizing (fits content instead of guessing a fixed size) --

_BULLET_SIZE_PRESETS = [
    (20, 16, 12, 6),
    (18, 15, 10, 5),
    (16, 13, 8, 4),
    (14.5, 12, 6, 3),
]


def _estimate_block_height_in(bullets, box_w_in, size0, size1, space0, space1, prefix0_len=3, prefix1_len=8):
    total = 0.0
    for level, text in bullets:
        size = size0 if level == 0 else size1
        space = space0 if level == 0 else space1
        prefix_len = prefix0_len if level == 0 else prefix1_len
        chars_per_line = max(10, int(box_w_in / (0.5 * size / 72)))
        lines = max(1, -(-(len(text) + prefix_len) // chars_per_line))
        line_height_in = (size * 1.22) / 72
        total += lines * line_height_in + space / 72
    return total


def _pick_bullet_sizes(bullets, box_w_in, box_h_in, presets=_BULLET_SIZE_PRESETS, safety_margin=0.88):
    # Require the estimate to clear a margin, not just technically fit --
    # the char-count-based line-wrap estimate is approximate, so an exact
    # fit on paper can still overflow once real font metrics render it.
    budget = box_h_in * safety_margin
    for preset in presets:
        if _estimate_block_height_in(bullets, box_w_in, *preset) <= budget:
            return preset
    return presets[-1]


def _fill_bullets(tf, bullets, size0, size1, space0, space1):
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space0 if level == 0 else space1)
        prefix = "▸  " if level == 0 else "     –  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size0 if level == 0 else size1)
        run.font.color.rgb = INK if level == 0 else MUTED_DARK
        run.font.name = FONT_BODY
        run.font.bold = level == 0


# ---- slide builders ---------------------------------------------------------

def add_title_slide(prs, kicker, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.5), kicker.upper(), 18, ACCENT, bold=True)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6), title, 44, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.35), Inches(11.3), Inches(1.0), subtitle, 20, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


def add_bullet_slide(prs, kicker, title, bullets, page, color=ACCENT, note=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.9), title, 30, INK, bold=True)

    box_w_in, box_h_in = 11.7, 4.9
    sizes = _pick_bullet_sizes(bullets, box_w_in, box_h_in)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(box_w_in), Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    _fill_bullets(tf, bullets, *sizes)

    if note:
        _textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45), note, 12, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_table_slide(prs, kicker, title, headers, rows, page, color=ACCENT, note=None):
    """A real small data table -- for showing actual numbers, not bullets pretending to be numbers."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.55 * n_rows)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = WHITE
                run.font.name = FONT_BODY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF2, 0xF4, 0xF8)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = INK
                    run.font.name = FONT_BODY

    if note:
        note_top = top + height + Inches(0.3)
        _textbox(slide, Inches(0.8), note_top, Inches(11.7), Inches(1.2), note, 14, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_closing_slide(prs, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2), title, 40, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(1.4), subtitle, 18, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


SERIES_MODULES = [
    (1, "OCR & Document Understanding", "PDF to structured markdown"),
    (2, "Analysis & Validation", "Is the extraction actually right?"),
    (3, "Chunking", "Splitting documents for retrieval"),
    (4, "Embedding & Database", "Vectors, Postgres, pgvector"),
    (5, "Database Deep Dive", "Search, indexes, and what Postgres actually does"),
    (6, "Summarization", "Retrieval + a real LLM + a UI"),
]


def add_series_intro_slide(prs, current_module, page):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT)
    _kicker(slide, "Video Series")
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(1.1), "Part of a Series: How to Build Your First Chatbot", 28, INK, bold=True)

    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "This module is one step in a planned video series that builds a real, working "
        "chatbot end to end -- real documents, real models, real measured results, not toy examples."
    )
    run.font.size = Pt(16)
    run.font.color.rgb = MUTED_DARK
    run.font.name = FONT_BODY
    run.font.italic = True

    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(11.7), Inches(3.0))
    tf2 = list_box.text_frame
    tf2.word_wrap = True
    first = True
    for num, name, desc in SERIES_MODULES:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(10)
        is_current = num == current_module
        run = p.add_run()
        prefix = "►  " if is_current else "    "
        suffix = "   ◂ you are here" if is_current else ""
        run.text = f"{prefix}Module {num}: {name} -- {desc}{suffix}"
        run.font.size = Pt(18) if is_current else Pt(16)
        run.font.bold = is_current
        run.font.color.rgb = ACCENT if is_current else MUTED_DARK
        run.font.name = FONT_BODY

    _textbox(
        slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
        "More modules are planned as the series continues toward a complete chatbot -- retrieval and response generation are next.",
        13, MUTED_DARK, italic=True,
    )

    _page_number(slide, page)
    return slide


# ---- deck content -----------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = [0]

    def next_page():
        page[0] += 1
        return page[0]

    # 1 -- title
    add_title_slide(
        prs,
        kicker="Module 2 · Analysis",
        title="Validating What OCR Actually Extracted",
        subtitle="Module 1 produces a markdown file and says it's done. This module asks: is it actually right?",
        footer="Identification  ·  Validation  ·  Root Cause Analysis",
    )
    next_page()

    # 2 -- series intro
    add_series_intro_slide(prs, current_module=2, page=next_page())

    # 3 -- why analysis
    add_bullet_slide(
        prs, "Why This Exists", "“Done” Is Not the Same as “Correct”",
        [
            (0, "Module 1 finishes every run the same way: a .md file and a metadata.json get written, no error, no warning"),
            (0, "That tells you the pipeline didn't crash. It tells you nothing about whether the output is trustworthy"),
            (0, "Module 2's job: turn the raw output of many documents into numbers you can actually look at and question"),
            (1, "page counts, token density, artifacts found, regions the model wasn't sure about, regions it dropped"),
            (0, "The deliverable is one Excel workbook -- because the person who needs to sanity-check 50 documents is rarely the person who wrote the pipeline"),
        ],
        next_page(),
    )

    # 4 -- how analysis helps
    add_bullet_slide(
        prs, "How This Helps", "Three Things This Module Actually Buys You",
        [
            (0, "Identification -- given 50 processed documents, which ones likely have an extraction problem, from the numbers alone, before opening a single file"),
            (0, "Validation -- confirm a document that finished without errors is actually trustworthy enough to move downstream into chunking and embedding"),
            (0, "Root cause analysis -- when a document does look wrong, trace it back to the exact page and region that caused it, not just know something's off"),
            (0, "All three come from the same source: real per-document metrics, computed once, put in front of a human who can actually question them"),
        ],
        next_page(),
    )

    # 5 -- overview
    add_bullet_slide(
        prs, "Overview", "What This Module Does",
        [
            (0, "Reads every completed document Module 1 produced -- the extracted text plus its layout metadata"),
            (0, "Computes, per document: page count, real token count, tokens/page, images and tables found, artifact counts, low-confidence region count, dropped region count"),
            (0, "Writes one Excel workbook: a Summary sheet (one row per document) plus detail sheets for anything worth a second look"),
        ],
        next_page(),
    )

    # 6 -- the core question
    add_bullet_slide(
        prs, "The Core Question", "How Many Tokens Per Page, Really?",
        [
            (0, "“Word count” is not “token count.” A word can be one token, several tokens, or a fraction of a shared token -- it depends on the tokenizer"),
            (0, "Why this matters here specifically: every downstream step (chunking, any embedding or LLM call) is budgeted in tokens, not words or characters"),
            (0, "So this module counts tokens with the SAME tokenizer the extraction model uses, rather than estimating -- an approximation here would quietly mislead every step after it"),
        ],
        next_page(),
    )

    # 7 -- real numbers table
    add_table_slide(
        prs, "Real Numbers", "Tokens Per Page, Three Real Documents",
        ["Document", "Pages", "Tokens", "Tokens / Page"],
        [
            ["federal_register_proclamation", 3, 1526, "508.7"],
            ["synthetic_sample", 1, 513, "513.0"],
            ["nasa_iss_factsheet", 2, 495, "247.5"],
        ],
        next_page(),
        note="The dense legal proclamation and the structured synthetic page (table + chart + formula) land almost identically dense -- ~510 tokens/page. The NASA fact sheet's plain prose runs at roughly half that. Content density drives this number more than document length does.",
    )

    # 8 -- layout analysis, the other axis
    add_bullet_slide(
        prs, "Layout Analysis", "The Other Half of Validation",
        [
            (0, "Token counting tells you how much content came out. It says nothing about whether the layout detector actually understood the page"),
            (0, "Module 1's layout detector already assigns a confidence score to every region it finds -- Module 2 doesn't compute this, it surfaces what already exists"),
            (1, "low_confidence_regions -- the detector found something there, but wasn't sure what"),
            (1, "dropped_regions -- the detector excluded a region entirely: correctly (headers/footers), or not (see the case study ahead)"),
            (0, "Same source also covers artifacts (logos, signatures, stamp/seals) and malformed tables -- structural failures, not content-volume ones"),
            (0, "Two independent questions: token analysis asks “how much came out,” layout analysis asks “did it even see the page correctly”"),
        ],
        next_page(),
    )

    # 9 -- excel report structure
    add_bullet_slide(
        prs, "The Deliverable", "One Workbook, Built for Validation",
        [
            (0, "Summary -- one row per document, every stat from the last two slides in one place, built to be scanned fast"),
            (0, "Artifacts Detail -- one row per detected logo/signature/stamp: page, description, exact location"),
            (0, "Low Confidence Regions -- one row per region the layout detector wasn't sure about"),
            (0, "Dropped Regions -- one row per region excluded from the document, including anything misclassified and dropped by mistake"),
            (1, "Every flagged row carries a location back to the original page -- a finding is only useful if it can be checked"),
        ],
        next_page(),
    )

    # 10 -- root cause analysis: real case study
    add_bullet_slide(
        prs, "Root Cause Analysis", "Tracing a Dropped Region Back to a Real Bug",
        [
            (0, "One document's Dropped Regions sheet flagged a region: low confidence, a specific location on the page"),
            (0, "That location turned out to be the document's logo -- a real image the layout detector misclassified as page furniture and dropped before extraction ever saw it"),
            (0, "Looking at that exact location on the original page image confirmed it immediately -- a visible box drawn right over the logo"),
            (1, "This is the whole point of keeping every dropped region instead of silently discarding them -- it turns a hidden bug into something traceable, and something fixable"),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 11 -- human in the loop
    add_bullet_slide(
        prs, "Human-in-the-Loop, By Design", "A Flag Is Only Useful If Someone Can Act On It",
        [
            (0, "Every flagged row -- low confidence, dropped region, malformed table -- carries a page number and a bounding box back to the original document"),
            (0, "That's not incidental: it's what turns a flag into a 30-second check instead of a full re-read of the document"),
            (0, "The loop: detector flags an uncertain region → it lands in the Excel workbook → a person looks at that exact spot on the original page → confirms a real problem, or clears it as a false alarm"),
            (1, "The logo case study on the previous slide is this loop working end to end -- a two-minute human check catching a bug automation alone missed"),
            (0, "What this module doesn't do yet: feed a human's correction back into re-thresholding or retraining the detector -- that's the natural next step, not built here"),
        ],
        next_page(),
    )

    # 12 -- validation as anomaly detection habit
    add_bullet_slide(
        prs, "A Validation Habit", "Tokens/Page as an Anomaly Signal",
        [
            (0, "There's no universal “correct” tokens-per-page number -- it depends entirely on the document type"),
            (0, "But once you have it for a batch of similar documents, an outlier is worth a look:"),
            (1, "Much LOWER than its peers -- possible sign of missed or truncated content"),
            (1, "Much HIGHER than its peers -- possible sign of garbled or repeated extraction output"),
            (0, "The value isn't any single number -- it's being able to scan 50 rows and immediately see which ones don't look like the rest"),
        ],
        next_page(),
    )

    # 13 -- sampling at scale
    add_bullet_slide(
        prs, "Sampling At Scale", "From 3 Documents to 3,000",
        [
            (0, "This module's corpus is 3 real documents -- reviewing every flagged row by hand is trivial at that size. It stops being trivial well before you reach a few thousand"),
            (0, "The fix isn't reviewing everything, and it isn't reviewing a random handful either -- it's spending human attention where the model already told you it's unsure"),
            (1, "Targeted review, 100% coverage -- every document with any low-confidence region, dropped region, or malformed table. Usually a small minority of the corpus, which is exactly why full coverage stays affordable"),
            (1, "Random audit sample of the rest -- a small percentage (5-10% is a common starting point) of documents that came back with zero flags, specifically to catch cases where the detector is confidently wrong"),
            (0, "Why the second bucket matters: a validation system that only ever double-checks its own doubts can't catch its own blind spots -- the audit sample is how you find out if “zero flags” is actually trustworthy"),
            (0, "Tokens/page outliers (previous slide) feed the same triage -- an outlier is one more reason to pull a document into the reviewed set even with zero layout flags"),
        ],
        next_page(),
        note="This is methodology, not a result -- honest about what a 3-document teaching corpus can and can't demonstrate. The same low_confidence_count / dropped_region_count / tokens_per_page columns already in the Excel report are exactly what a sampling rule would key off of at real scale.",
    )

    # 14 -- what counts as a document
    add_bullet_slide(
        prs, "A Quiet Detail", "Partial Runs Get Skipped, Not Guessed At",
        [
            (0, "A document only gets analyzed once its extraction has actually finished, not while it's still in progress"),
            (0, "An interrupted or still-running extraction leaves incomplete output"),
            (0, "That incomplete output gets skipped rather than guessed at -- a document either finished cleanly, or it isn't validated at all"),
        ],
        next_page(),
    )

    # 15 -- wrap-up
    add_bullet_slide(
        prs, "Wrap-Up", "Why This Module Matters",
        [
            (0, "It turns “the pipeline didn't crash” into “here's proof this document is actually usable”"),
            (0, "Two independent checks, not one -- token analysis for how much came out, layout analysis for whether the model saw the page correctly"),
            (0, "It's human-in-the-loop by design -- every flag carries an exact way back to the page a person needs to look at"),
            (0, "It scales past hand-checking every document -- targeted review of what's flagged, random audits of what isn't, once the corpus is bigger than 3 files"),
            (0, "Next: Module 3 takes this same validated output and turns it into retrieval-ready chunks -- recursive, semantic, and hierarchical strategies, compared on the same real documents"),
        ],
        next_page(),
    )

    # 17 -- closing
    add_closing_slide(
        prs,
        title="Questions?",
        subtitle="2. ANALYSIS/  --  README.md has the full column-by-column report reference and every real number in this deck.",
        footer="Module 2 of 3  ·  Data Processing → Analysis → Chunking",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
