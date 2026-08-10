"""Builds Database_Module_Overview.pptx -- a teaching deck for Module 5 (Database).

Every number in this deck comes from actually running run_demos.py against
Module 4's live Postgres/pgvector container -- not illustrative placeholders.
Re-run run_demos.py first if you want the deck to reflect a fresh run; this
script doesn't re-run it for you.

Same visual design system as the other four modules' decks, duplicated here
rather than imported -- each module stays self-contained and only shares
the venv.

Usage:
  python make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "Database_Module_Overview.pptx"

# ---- palette (matches the other four modules' decks) -------------------------
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x27)
MUTED = RGBColor(0x8A, 0x93, 0xA6)
MUTED_DARK = RGBColor(0x5B, 0x63, 0x74)
ACCENT = RGBColor(0x3E, 0x8E, 0xDE)       # blue -- structural/technical content
ACCENT_WARM = RGBColor(0xF2, 0xA6, 0x3D)  # amber -- "did you know" slides
ACCENT_GOOD = RGBColor(0x4C, 0xAF, 0x7D)  # green -- real code slides
ACCENT_BAD = RGBColor(0xE0, 0x6C, 0x5C)   # red -- real findings / case studies
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT = RGBColor(0xA6, 0xE3, 0xA1)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def _kicker(slide, text, color=ACCENT):
    _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5), text.upper(), 15, color, bold=True)


def _accent_bar(slide, color, top=Inches(0), height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _page_number(slide, n, dark=False):
    color = MUTED if dark else MUTED_DARK
    _textbox(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35), str(n), 11, color, align=PP_ALIGN.RIGHT)


# ---- adaptive bullet sizing (fits content instead of guessing a fixed size) --

_BULLET_SIZE_PRESETS = [
    (20, 16, 12, 6),
    (18, 15, 10, 5),
    (16, 13, 8, 4),
    (14.5, 12, 6, 3),
]


def _estimate_block_height_in(bullets, box_w_in, size0, size1, space0, space1, prefix0_len=3, prefix1_len=8):
    total = 0.0
    for level, text in bullets:
        size = size0 if level == 0 else size1
        space = space0 if level == 0 else space1
        prefix_len = prefix0_len if level == 0 else prefix1_len
        chars_per_line = max(10, int(box_w_in / (0.5 * size / 72)))
        lines = max(1, -(-(len(text) + prefix_len) // chars_per_line))
        line_height_in = (size * 1.22) / 72
        total += lines * line_height_in + space / 72
    return total


def _pick_bullet_sizes(bullets, box_w_in, box_h_in, presets=_BULLET_SIZE_PRESETS, safety_margin=0.88):
    budget = box_h_in * safety_margin
    for preset in presets:
        if _estimate_block_height_in(bullets, box_w_in, *preset) <= budget:
            return preset
    return presets[-1]


def _fill_bullets(tf, bullets, size0, size1, space0, space1):
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space0 if level == 0 else space1)
        prefix = "▸  " if level == 0 else "     –  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size0 if level == 0 else size1)
        run.font.color.rgb = INK if level == 0 else MUTED_DARK
        run.font.name = FONT_BODY
        run.font.bold = level == 0


# ---- slide builders ---------------------------------------------------------

def add_title_slide(prs, kicker, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.5), kicker.upper(), 18, ACCENT, bold=True)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6), title, 44, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.35), Inches(11.3), Inches(1.0), subtitle, 20, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


# Real numbers from an actual run_demos.py run against a 200,000-row synthetic
# benchmark table -- see db_features/index_benchmark.py and the README.
INDEX_BENCHMARK_ROWS = [
    ["No index", "Gather Merge (parallel seq scan)", "23.18", "--"],
    ["HNSW", "Index Scan", "1.58", "215.6 s"],
    ["IVFFlat (lists=447)", "Index Scan", "0.84", "7.8 s"],
]

SERIES_MODULES = [
    (1, "OCR & Document Understanding", "PDF to structured markdown"),
    (2, "Analysis & Validation", "Is the extraction actually right?"),
    (3, "Chunking", "Splitting documents for retrieval"),
    (4, "Embedding & Database", "Vectors, Postgres, pgvector"),
    (5, "Database Deep Dive", "Search, indexes, and what Postgres actually does"),
    (6, "Summarization", "Retrieval + a real LLM + a UI"),
]


def add_series_intro_slide(prs, current_module, page):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT)
    _kicker(slide, "Video Series")
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(1.1), "Part of a Series: How to Build Your First Chatbot", 28, INK, bold=True)

    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "This module is one step in a planned video series that builds a real, working "
        "chatbot end to end -- real documents, real models, real measured results, not toy examples."
    )
    run.font.size = Pt(16)
    run.font.color.rgb = MUTED_DARK
    run.font.name = FONT_BODY
    run.font.italic = True

    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(3.1))
    tf2 = list_box.text_frame
    tf2.word_wrap = True
    first = True
    for num, name, desc in SERIES_MODULES:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(9)
        is_current = num == current_module
        run = p.add_run()
        prefix = "►  " if is_current else "    "
        suffix = "   ◂ you are here" if is_current else ""
        run.text = f"{prefix}Module {num}: {name} -- {desc}{suffix}"
        run.font.size = Pt(17) if is_current else Pt(15)
        run.font.bold = is_current
        run.font.color.rgb = ACCENT if is_current else MUTED_DARK
        run.font.name = FONT_BODY

    _textbox(
        slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.8),
        "More modules are planned as the series continues toward a complete chatbot -- retrieval and response generation are next.",
        13, MUTED_DARK, italic=True,
    )

    _page_number(slide, page)
    return slide


def add_bullet_slide(prs, kicker, title, bullets, page, color=ACCENT, note=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.9), title, 30, INK, bold=True)

    box_w_in, box_h_in = 11.7, 4.9
    sizes = _pick_bullet_sizes(bullets, box_w_in, box_h_in)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(box_w_in), Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    _fill_bullets(tf, bullets, *sizes)

    if note:
        _textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45), note, 12, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_fact_slide(prs, kicker, title, bullets, page, color=ACCENT_WARM):
    slide = add_bullet_slide(prs, kicker, title, bullets, page, color=color)
    tag = slide.shapes.add_textbox(Inches(10.6), Inches(0.42), Inches(2.0), Inches(0.5))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "★ DID YOU KNOW"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_WARM
    return slide


def add_code_slide(prs, kicker, title, label, code_text, note, page, height_in=4.0):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT_GOOD)
    _kicker(slide, kicker, color=ACCENT_GOOD)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)
    _textbox(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4), label, 14, MUTED_DARK, bold=True)

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.7), Inches(height_in))
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line if line.strip() else " "
        run.font.size = Pt(14)
        run.font.name = FONT_CODE
        run.font.color.rgb = CODE_TEXT

    note_top = 2.0 + height_in + 0.2
    _textbox(slide, Inches(0.8), Inches(note_top), Inches(11.7), Inches(7.0 - note_top), note, 14, MUTED_DARK, italic=True)
    _page_number(slide, page)
    return slide


def add_table_slide(prs, kicker, title, headers, rows, page, color=ACCENT, note=None, col_widths=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.8), Inches(1.95), Inches(11.7), Inches(0.5 * n_rows)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table

    if col_widths:
        for c, w in enumerate(col_widths):
            table.columns[c].width = Inches(w)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = WHITE
                run.font.name = FONT_BODY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF2, 0xF4, 0xF8)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = INK
                    run.font.name = FONT_BODY

    if note:
        note_top = top + height + Inches(0.25)
        _textbox(slide, Inches(0.8), note_top, Inches(11.7), Inches(7.2) - note_top, note, 13, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_closing_slide(prs, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2), title, 40, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(1.4), subtitle, 18, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


# ---- deck content -----------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = [0]

    def next_page():
        page[0] += 1
        return page[0]

    # 1 -- title
    add_title_slide(
        prs,
        kicker="Module 5 · Database",
        title="A Deep Dive Into Postgres + pgvector",
        subtitle="Not a new database -- the one Module 4 already built. Distance operators, metadata filtering, full-text search, and what actually happens when you add an index.",
        footer="Postgres 16  ·  pgvector  ·  tsvector/GIN  ·  HNSW  ·  IVFFlat",
    )
    next_page()

    # 2 -- series intro
    add_series_intro_slide(prs, current_module=5, page=next_page())

    # 3 -- fun fact: postgres history
    add_fact_slide(
        prs, "Did you know", "Postgres Predates the Web",
        [
            (0, "Postgres traces back to POSTGRES at UC Berkeley in 1986 -- led by Michael Stonebraker, who later won the Turing Award largely for this work"),
            (0, "It's been continuously developed for close to 40 years, longer than almost any other actively used production database"),
            (0, "\"Just use Postgres\" became a real engineering meme for a reason: extensions like pgvector let it absorb entirely new categories of workload (vector search, in this case) without leaving the database at all"),
        ],
        next_page(),
    )

    # 4 -- why postgres for this stack
    add_bullet_slide(
        prs, "Why Postgres, Specifically", "One Database, Not a Pile of Specialized Ones",
        [
            (0, "Extensions, not a separate product: CREATE EXTENSION vector; -- vector search lives in the same database as everything else, one connection, one transaction model"),
            (0, "Real SQL around the vectors: WHERE doc_name = ... combined with ORDER BY embedding <=> ... is one query, not \"search a vector store, then filter in application code\""),
            (0, "More than one kind of search, natively: full-text search (tsvector/tsquery) ships with core Postgres, no extension needed"),
            (0, "A real, inspectable query planner: EXPLAIN ANALYZE tells you the truth about whether an index is actually used, not just whether you created one"),
        ],
        next_page(),
    )

    # 5 -- distance operators
    add_code_slide(
        prs, "pgvector Feature 1", "Three Distance Operators",
        "pgvector operators -- <-> Euclidean, <=> cosine, <#> negative inner product",
        "SELECT chunk_id, embedding_256 <=> %s AS distance\n"
        "FROM chunks\n"
        "ORDER BY embedding_256 <=> %s\n"
        "LIMIT 5;",
        "Tutorials often say \"use cosine for normalized embeddings\" as though the choice is consequential. The next slide checks that empirically instead of repeating it.",
        next_page(),
        height_in=2.2,
    )

    # 6 -- real finding: operators identical ranking
    add_bullet_slide(
        prs, "Real Finding", "All Three Operators Agree -- and Here's Why",
        [
            (0, "Module 4's embeddings are L2-normalized (unit vectors). For unit vectors, L2² = 2 − 2·cos_similarity -- a strictly monotonic function of cosine similarity"),
            (0, "So <->, <=>, and <#> should produce the IDENTICAL ranking on this data, even though their raw numbers look completely different"),
            (0, "Tested on 3 real queries against the real chunks table: confirmed identical top-5 rankings on every single one -- verified, not assumed"),
            (1, "The raw distance values differ (0.26 vs 0.44 vs -0.65 for the same pair, on the same query) -- but WHO WINS never changes. That's the number that actually matters for retrieval."),
        ],
        next_page(),
    )

    # 7 -- fun fact: vector math
    add_fact_slide(
        prs, "Did you know", "\"Cosine Similarity\" Is Just the Angle Between Two Arrows",
        [
            (0, "Cosine similarity is literally the cosine of the angle between two vectors -- 1.0 means pointing the same direction, 0 means perpendicular (unrelated), -1 means opposite"),
            (0, "It ignores vector LENGTH entirely, which is exactly why embeddings get L2-normalized before comparison -- you want \"same meaning, different confidence\" to still count as similar"),
            (0, "This is 19th-century vector algebra doing the heavy lifting inside a 2020s AI pipeline -- nothing about the comparison itself is new"),
        ],
        next_page(),
    )

    # 8 -- metadata filtering
    add_code_slide(
        prs, "pgvector Feature 2", "Metadata-Filtered Vector Search",
        "One query: a WHERE filter AND a vector ORDER BY, together",
        "SELECT chunk_id, chunk_text\n"
        "FROM chunks\n"
        "WHERE doc_name = 'nasa_iss_factsheet'\n"
        "ORDER BY embedding_256 <=> %s\n"
        "LIMIT 5;",
        "\"Find content similar to X, but only within document Y\" -- the real pattern retrieval systems actually use, not a special pgvector syntax.",
        next_page(),
        height_in=2.4,
    )

    # 9 -- real finding: filter changes result
    add_bullet_slide(
        prs, "Real Finding", "A Case Where Filtering Actually Changes the Answer",
        [
            (0, "Query: \"important developments and changes\""),
            (0, "Unfiltered: top result comes from the Federal Register document (a chunk about manufacturing and tax policy)"),
            (0, "Filtered to doc_name = 'nasa_iss_factsheet': top result correctly becomes the NASA chunk about Space Station research"),
            (1, "Tested several other queries first -- most didn't change, since a 3-document, topically distinct corpus doesn't confuse vector search easily. This one genuinely does."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 10 -- fun fact: full text search history
    add_fact_slide(
        prs, "Did you know", "Postgres Has Had Real Full-Text Search Since 2008",
        [
            (0, "tsvector/tsquery (Postgres's built-in text search) shipped in Postgres 8.3, in 2008 -- long before \"add search to your app\" became a SaaS product category"),
            (0, "It supports stemming (\"running\" matches \"run\"), ranking (ts_rank), and GIN indexes for speed -- a genuinely complete search engine, not a toy LIKE '%...%' substitute"),
            (0, "Most people reach for a separate search service by default without checking whether Postgres already does what they need"),
        ],
        next_page(),
    )

    # 11 -- fulltext vs vector table
    add_table_slide(
        prs, "Real Numbers", "Full-Text vs. Vector Search, Same Queries",
        ["Query", "Full-text result", "Vector result"],
        [
            ['"92 percent"', "correct, rank 1", "WRONG at rank 1 (correct is rank 2)"],
            ['"Houston Texas"', "correct", "correct"],
            ['"orbital laboratory"', "correct", "correct"],
            ['"how much did border crossings decrease" (paraphrase)', "no matches at all", "correct"],
            ['"an orbiting laboratory where things float" (paraphrase)', "no matches at all", "correct"],
        ],
        next_page(),
        note="Full-text wins on exact terms and numbers. Vector wins on paraphrase. Neither replaces the other -- this is the actual case for hybrid search, not a hypothetical one.",
        col_widths=[4.8, 3.3, 3.6],
    )

    # 12 -- index types concept
    add_bullet_slide(
        prs, "pgvector Feature 3", "Two Kinds of Approximate Index",
        [
            (0, "Both trade a small amount of accuracy for a large amount of speed on big tables -- neither is exact nearest-neighbor search, by design"),
            (0, "IVFFlat: clusters vectors into \"lists\" ahead of time, then only searches the nearest clusters at query time. Fast to build, needs a lists parameter tuned to table size"),
            (0, "HNSW: builds a navigable graph connecting similar vectors. Slower to build, but typically better accuracy at query time, and needs no tuning parameter to get started"),
            (0, "Neither one matters on a small table -- Postgres's planner will (correctly) ignore both and just scan everything. The next slides prove that with real numbers."),
        ],
        next_page(),
    )

    # 13 -- real code: index creation
    add_code_slide(
        prs, "The Actual Code", "Building Both Index Types",
        "db_features/index_benchmark.py",
        "CREATE INDEX ... ON benchmark_vectors\n"
        "  USING hnsw (embedding vector_cosine_ops);\n\n"
        "CREATE INDEX ... ON benchmark_vectors\n"
        "  USING ivfflat (embedding vector_cosine_ops)\n"
        "  WITH (lists = 447);   -- pgvector's own rule of thumb: sqrt(row_count)",
        "vector_cosine_ops tells the index to optimize specifically for the <=> operator -- an index built for one distance operator doesn't necessarily help another.",
        next_page(),
        height_in=2.9,
    )

    # 14 -- real numbers: index benchmark
    add_table_slide(
        prs, "Real Numbers", "200,000 Synthetic Vectors: No Index vs. HNSW vs. IVFFlat",
        ["Method", "Scan Type (EXPLAIN ANALYZE)", "Wall-clock (ms)", "Index build time"],
        INDEX_BENCHMARK_ROWS,
        next_page(),
        note="Table is synthetic -- 200,000 random unit vectors, built specifically to have enough rows for an index to matter. Module 4's real 18-chunk table is too small for any index to change anything.",
        col_widths=[2.2, 4.3, 2.6, 2.6],
    )

    # 15 -- real finding: build time tradeoff
    add_fact_slide(
        prs, "Real Finding", "215 Seconds vs. 7.8 Seconds to Build",
        [
            (0, "Both indexes made queries dramatically faster: 23.18 ms with no index down to 1.58 ms (HNSW) and 0.84 ms (IVFFlat) -- roughly 15x and 28x"),
            (0, "But building the HNSW index took 215.6 seconds on 200,000 rows. IVFFlat took 7.8 seconds -- about 27x faster to build"),
            (0, "IVFFlat also edged out HNSW on query speed here, which isn't the usual story -- HNSW is normally the accuracy/recall favorite. Likely explanation: these vectors are uniformly random with no real semantic cluster structure, which is exactly what IVFFlat's clustering approach is suited for and doesn't showcase HNSW's usual advantage"),
            (1, "The lesson isn't \"IVFFlat wins\" -- it's that the right index depends on your actual data and your actual constraints (build time vs. query time vs. accuracy), not a universal default."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 16 -- real finding: explain analyze tells the truth
    add_bullet_slide(
        prs, "Real Finding", "Creating an Index Isn't the Same as Using One",
        [
            (0, "CREATE INDEX doesn't guarantee Postgres will use it -- the query planner decides, based on real cost estimates, at query time"),
            (0, "On Module 4's 18-row chunks table, an index would sit there unused; the planner correctly prefers a sequential scan because the table is too small for the index to pay for itself"),
            (0, "EXPLAIN (ANALYZE, FORMAT JSON) is how you check the truth instead of assuming -- it names the actual scan node Postgres chose, every time"),
            (1, "\"I created an index\" is a claim. \"EXPLAIN ANALYZE shows an Index Scan\" is evidence."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 17 -- wrap-up
    add_bullet_slide(
        prs, "Wrap-Up", "What Module 5 Actually Covered",
        [
            (0, "Three distance operators, empirically confirmed to rank identically on normalized embeddings -- and why that's true, not just observed"),
            (0, "Metadata-filtered vector search, with a real query where the filter changes the answer"),
            (0, "Full-text search as a genuine complement to vector search, not a lesser substitute -- each wins on exactly the cases you'd expect"),
            (0, "A real, measured index benchmark -- and a reminder that EXPLAIN ANALYZE, not intuition, is how you know an index is actually helping"),
        ],
        next_page(),
    )

    # 18 -- closing
    add_closing_slide(
        prs,
        title="Questions?",
        subtitle="5. DATABASE/  --  README.md has the full setup and every real number in this deck.",
        footer="Module 5  ·  Data Processing → Analysis → Chunking → Embedding & Database → Database Deep Dive",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
