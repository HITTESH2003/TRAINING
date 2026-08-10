"""Builds Summarisation_Module_Overview.pptx -- a teaching deck for
Module 6 (Summarization).

Every number in this deck comes from actually running run_topk_experiment.py
against Module 4's live Postgres/pgvector container and a real Qwen3-0.6B
generation -- not illustrative placeholders. Re-run it first if you want
the deck to reflect a fresh run; this script doesn't re-run it for you.

Same visual design system as the other five modules' decks, duplicated
here rather than imported -- each module stays self-contained and only
shares the venv.

Usage:
  python make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "Summarisation_Module_Overview.pptx"

# ---- palette (matches the other five modules' decks) -------------------------
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x27)
MUTED = RGBColor(0x8A, 0x93, 0xA6)
MUTED_DARK = RGBColor(0x5B, 0x63, 0x74)
ACCENT = RGBColor(0x3E, 0x8E, 0xDE)       # blue -- structural/technical content
ACCENT_WARM = RGBColor(0xF2, 0xA6, 0x3D)  # amber -- "did you know" slides
ACCENT_GOOD = RGBColor(0x4C, 0xAF, 0x7D)  # green -- real code slides
ACCENT_BAD = RGBColor(0xE0, 0x6C, 0x5C)   # red -- real findings / case studies
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT = RGBColor(0xA6, 0xE3, 0xA1)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def _kicker(slide, text, color=ACCENT):
    _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5), text.upper(), 15, color, bold=True)


def _accent_bar(slide, color, top=Inches(0), height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _page_number(slide, n, dark=False):
    color = MUTED if dark else MUTED_DARK
    _textbox(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35), str(n), 11, color, align=PP_ALIGN.RIGHT)


# ---- adaptive bullet sizing (fits content instead of guessing a fixed size) --

_BULLET_SIZE_PRESETS = [
    (20, 16, 12, 6),
    (18, 15, 10, 5),
    (16, 13, 8, 4),
    (14.5, 12, 6, 3),
]


def _estimate_block_height_in(bullets, box_w_in, size0, size1, space0, space1, prefix0_len=3, prefix1_len=8):
    total = 0.0
    for level, text in bullets:
        size = size0 if level == 0 else size1
        space = space0 if level == 0 else space1
        prefix_len = prefix0_len if level == 0 else prefix1_len
        chars_per_line = max(10, int(box_w_in / (0.5 * size / 72)))
        lines = max(1, -(-(len(text) + prefix_len) // chars_per_line))
        line_height_in = (size * 1.22) / 72
        total += lines * line_height_in + space / 72
    return total


def _pick_bullet_sizes(bullets, box_w_in, box_h_in, presets=_BULLET_SIZE_PRESETS, safety_margin=0.88):
    budget = box_h_in * safety_margin
    for preset in presets:
        if _estimate_block_height_in(bullets, box_w_in, *preset) <= budget:
            return preset
    return presets[-1]


def _fill_bullets(tf, bullets, size0, size1, space0, space1):
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space0 if level == 0 else space1)
        prefix = "▸  " if level == 0 else "     –  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size0 if level == 0 else size1)
        run.font.color.rgb = INK if level == 0 else MUTED_DARK
        run.font.name = FONT_BODY
        run.font.bold = level == 0


# ---- slide builders ---------------------------------------------------------

def add_title_slide(prs, kicker, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.5), kicker.upper(), 18, ACCENT, bold=True)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6), title, 44, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.35), Inches(11.3), Inches(1.0), subtitle, 20, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


SERIES_MODULES = [
    (1, "OCR & Document Understanding", "PDF to structured markdown"),
    (2, "Analysis & Validation", "Is the extraction actually right?"),
    (3, "Chunking", "Splitting documents for retrieval"),
    (4, "Embedding & Database", "Vectors, Postgres, pgvector"),
    (5, "Database Deep Dive", "Search, indexes, and what Postgres actually does"),
    (6, "Summarization", "Retrieval + a real LLM + a UI"),
]


def add_series_intro_slide(prs, current_module, page):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT)
    _kicker(slide, "Video Series")
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(1.1), "Part of a Series: How to Build Your First Chatbot", 28, INK, bold=True)

    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "This module is one step in a planned video series that builds a real, working "
        "chatbot end to end -- real documents, real models, real measured results, not toy examples."
    )
    run.font.size = Pt(16)
    run.font.color.rgb = MUTED_DARK
    run.font.name = FONT_BODY
    run.font.italic = True

    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.95), Inches(11.7), Inches(3.35))
    tf2 = list_box.text_frame
    tf2.word_wrap = True
    first = True
    for num, name, desc in SERIES_MODULES:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(8)
        is_current = num == current_module
        run = p.add_run()
        prefix = "►  " if is_current else "    "
        suffix = "   ◂ you are here" if is_current else ""
        run.text = f"{prefix}Module {num}: {name} -- {desc}{suffix}"
        run.font.size = Pt(16) if is_current else Pt(14)
        run.font.bold = is_current
        run.font.color.rgb = ACCENT if is_current else MUTED_DARK
        run.font.name = FONT_BODY

    _textbox(
        slide, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.7),
        "This module is the payoff -- everything upstream (extraction, validation, chunking, embedding, storage) finally answers a real question.",
        13, MUTED_DARK, italic=True,
    )

    _page_number(slide, page)
    return slide


def add_bullet_slide(prs, kicker, title, bullets, page, color=ACCENT, note=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.9), title, 30, INK, bold=True)

    box_w_in, box_h_in = 11.7, 4.9
    sizes = _pick_bullet_sizes(bullets, box_w_in, box_h_in)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(box_w_in), Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    _fill_bullets(tf, bullets, *sizes)

    if note:
        _textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45), note, 12, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_fact_slide(prs, kicker, title, bullets, page, color=ACCENT_WARM):
    slide = add_bullet_slide(prs, kicker, title, bullets, page, color=color)
    tag = slide.shapes.add_textbox(Inches(10.6), Inches(0.42), Inches(2.0), Inches(0.5))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "★ DID YOU KNOW"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_WARM
    return slide


def add_code_slide(prs, kicker, title, label, code_text, note, page, height_in=4.0):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT_GOOD)
    _kicker(slide, kicker, color=ACCENT_GOOD)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)
    _textbox(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4), label, 14, MUTED_DARK, bold=True)

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.7), Inches(height_in))
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line if line.strip() else " "
        run.font.size = Pt(14)
        run.font.name = FONT_CODE
        run.font.color.rgb = CODE_TEXT

    note_top = 2.0 + height_in + 0.2
    _textbox(slide, Inches(0.8), Inches(note_top), Inches(11.7), Inches(7.0 - note_top), note, 14, MUTED_DARK, italic=True)
    _page_number(slide, page)
    return slide


def add_table_slide(prs, kicker, title, headers, rows, page, color=ACCENT, note=None, col_widths=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.8), Inches(1.95), Inches(11.7), Inches(0.5 * n_rows)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table

    if col_widths:
        for c, w in enumerate(col_widths):
            table.columns[c].width = Inches(w)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = WHITE
                run.font.name = FONT_BODY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF2, 0xF4, 0xF8)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = INK
                    run.font.name = FONT_BODY

    if note:
        note_top = top + height + Inches(0.25)
        _textbox(slide, Inches(0.8), note_top, Inches(11.7), Inches(7.2) - note_top, note, 13, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_closing_slide(prs, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2), title, 40, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(1.4), subtitle, 18, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


# ---- deck content -----------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = [0]

    def next_page():
        page[0] += 1
        return page[0]

    # 1 -- title
    add_title_slide(
        prs,
        kicker="Module 6 · Summarization",
        title="Retrieval + a Real LLM + a UI",
        subtitle="A question goes in, Module 4's database gets searched, Qwen3-0.6B turns the results into an answer -- every parameter along the way is a visible control, not a hidden constant.",
        footer="Qwen/Qwen3-0.6B  ·  Streamlit  ·  Postgres + pgvector",
    )
    next_page()

    # 2 -- series intro
    add_series_intro_slide(prs, current_module=6, page=next_page())

    # 3 -- why this module
    add_bullet_slide(
        prs, "Why This Exists", "Everything Upstream Finally Answers a Question",
        [
            (0, "Module 1 extracted. Module 2 validated. Module 3 chunked. Modules 4 and 5 embedded, stored, and searched."),
            (0, "None of that is useful to an actual user until something turns \"here are 5 relevant chunks\" into \"here's your answer\""),
            (0, "That's this module: real retrieval from a real database, real generation from a real small LLM, wrapped in a UI where every parameter is a slider you can actually move"),
        ],
        next_page(),
    )

    # 4 -- what is an LLM
    add_bullet_slide(
        prs, "Fundamentals", "What Is an LLM, Actually?",
        [
            (0, "A neural network trained to predict the next word (token) in a sequence, given everything before it -- nothing more mysterious than that at the core"),
            (0, "Trained on an enormous amount of text -- trillions of tokens -- until \"predict the next word\" produces fluent, coherent, contextually appropriate output"),
            (0, "It doesn't look facts up anywhere. Everything it \"knows\" is compressed into its parameters (weights) during training -- there's no database inside it"),
            (0, "Given a prompt, it generates a response one token at a time, each new token chosen based on everything generated so far"),
        ],
        next_page(),
    )

    # 5 -- what the llm does in this pipeline
    add_bullet_slide(
        prs, "Fundamentals", "What the LLM Actually Does Here",
        [
            (0, "In a RAG (Retrieval-Augmented Generation) system, the LLM is NOT the source of truth -- the retrieved chunks are"),
            (0, "Its job is narrower and more reliable: read the provided source excerpts, and write a coherent, grounded answer using ONLY that context"),
            (0, "That's a reading-comprehension-and-synthesis task, not \"recall a fact from memory\" -- a much easier, more constrained job than open-domain question answering"),
            (1, "This distinction is the whole reason a small model can do this job well, even though it couldn't reliably answer the same question from memory alone."),
        ],
        next_page(),
    )

    # 6 -- how the pipeline works
    add_bullet_slide(
        prs, "Fundamentals", "How the Pipeline Actually Works",
        [
            (0, "1. User asks a question"),
            (0, "2. The question gets embedded -- Nomic, the same embedding model from Module 4"),
            (0, "3. Postgres/pgvector searches for the top_k most similar chunks"),
            (0, "4. Retrieved chunks get assembled into a prompt, numbered and source-attributed"),
            (0, "5. Qwen3-0.6B reads that prompt and generates an answer"),
            (0, "6. The answer, plus which sources it came from, gets shown back to the user"),
            (1, "Retrieval and generation are two completely separate models doing two completely separate jobs -- one finds, one writes."),
        ],
        next_page(),
    )

    # 7 -- why 0.6B
    add_bullet_slide(
        prs, "Fundamentals", "Why a 0.6B Model, Specifically",
        [
            (0, "The LLM's job here is narrow -- read provided context, write a grounded answer -- not recall obscure facts from memory"),
            (0, "A narrow job doesn't need a 70B-parameter model: Qwen3-0.6B runs in a few seconds on a laptop's own GPU, no API key, no per-token cost, no network round trip"),
            (0, "Teaching-friendly on purpose: small enough to actually run live in a classroom, fast enough to change a parameter and immediately see the effect"),
            (1, "Real tradeoff, not hidden: this module's own experiment catches it making real mistakes. Small and free doesn't mean flawless -- see the findings a few slides from now."),
        ],
        next_page(),
    )

    # 8 -- model parameters explained
    add_bullet_slide(
        prs, "Fundamentals", "The Generation Parameters, Explained",
        [
            (0, "temperature -- how much randomness in word choice. Low = safe and predictable, high = varied and less predictable"),
            (0, "top_p (nucleus sampling) -- only sample from the smallest set of next-words whose combined probability reaches p, narrowing the field before picking"),
            (0, "top_k -- only consider the k most likely next words at each step, full stop"),
            (0, "max_new_tokens -- a hard cap on how long the generated answer is allowed to get"),
            (0, "enable_thinking -- whether the model writes out reasoning before its final answer (covered in detail next)"),
        ],
        next_page(),
    )

    # 9 -- fun fact: small llm capability
    add_fact_slide(
        prs, "Did you know", "0.6 Billion Parameters Is Genuinely Tiny for an LLM",
        [
            (0, "Qwen3-0.6B is roughly 300x smaller than GPT-3 (175B) and still small next to most \"small\" open models people mean when they say that (7-8B)"),
            (0, "It's small enough to run real generation on a laptop's own GPU in a few seconds -- no API key, no per-token bill, no network round trip"),
            (0, "The tradeoff is real too, not hidden: this module's own experiment catches it making real mistakes. Small and free doesn't mean flawless -- see the findings a few slides from now"),
        ],
        next_page(),
    )

    # 10 -- model config: enable_thinking
    add_bullet_slide(
        prs, "Model Config 1 of 2", "enable_thinking Isn't a Post-Processing Switch",
        [
            (0, "Qwen3 models can emit a <think>...</think> reasoning block before the actual answer -- but it's controlled when the PROMPT is built, not after generation"),
            (0, "apply_chat_template(messages, enable_thinking=True/False) changes what the model is asked to produce, not just how the output gets parsed"),
            (0, "Thinking content is split from the answer using the model card's own method: find token id 151668 (</think>) in the generated ids and split there -- not a string search on decoded text, which can be fooled by the model discussing \"</think>\" as a literal string"),
        ],
        next_page(),
    )

    # 11 -- model config: sampling presets
    add_table_slide(
        prs, "Model Config 2 of 2", "Sampling Parameters Differ by Mode, Officially",
        ["Mode", "Temperature", "top_p", "top_k"],
        [
            ["Thinking (enable_thinking=True)", "0.6", "0.95", "20"],
            ["Non-thinking (enable_thinking=False)", "0.7", "0.8", "20"],
        ],
        next_page(),
        note="Straight from Qwen3's own model card -- not one generic temperature reused for both modes. rag/config.py::GENERATION_PRESETS encodes exactly this.",
        col_widths=[5.5, 2.5, 2.0, 1.7],
    )

    # 12 -- real code: thinking split
    add_code_slide(
        prs, "The Actual Code", "Splitting Thinking From the Answer",
        "rag/llm.py -- the model card's own documented method",
        "try:\n"
        "    split = len(output_ids) - output_ids[::-1].index(151668)\n"
        "    thinking_ids, answer_ids = output_ids[:split], output_ids[split:]\n"
        "except ValueError:\n"
        "    thinking_ids, answer_ids = [], output_ids   # no </think> emitted",
        "151668 is the </think> token id. Searching from the end handles the model discussing thinking/reasoning inside its own answer without false-splitting.",
        next_page(),
        height_in=2.2,
    )

    # 13 -- the four parameters
    add_bullet_slide(
        prs, "What This Module Teaches", "Four Parameters, All Real Controls",
        [
            (0, "top_k -- how many retrieved chunks go into the prompt"),
            (0, "Embedding dimension -- the same Matryoshka truncation from Module 4, reused here"),
            (0, "enable_thinking -- changes both the sampling preset AND whether a reasoning trace gets generated"),
            (0, "max_new_tokens -- the hard cap on how long an answer can get"),
            (1, "Every one of these is a slider or toggle in the Streamlit app -- not a constant buried in code you'd need to edit to test."),
        ],
        next_page(),
    )

    # 14 -- the experiment
    add_bullet_slide(
        prs, "The Experiment", "One Question, top_k Swept From 1 to the Whole Corpus",
        [
            (0, "Fixed question: \"What does this collection of documents cover?\""),
            (0, "top_k tested: 1, 3, 6, 10, and 18 (the entire real 3-document, 18-chunk corpus from Modules 1-4)"),
            (0, "For each: which documents got retrieved, how many prompt tokens that cost, the actual generated answer, and real wall-clock generation time"),
        ],
        next_page(),
    )

    # 15 -- real numbers table
    add_table_slide(
        prs, "Real Numbers", "top_k vs. Document Coverage",
        ["top_k", "Documents represented", "Prompt tokens", "Generation time"],
        [
            ["1", "federal_register only", "245", "2.9 s"],
            ["3", "federal_register, nasa", "516", "2.2 s"],
            ["6", "federal_register, nasa", "1044", "3.6 s"],
            ["10", "ALL 3", "1560", "51.4 s"],
            ["18 (everything)", "all 3", "2712", "7.4 s"],
        ],
        next_page(),
        note="Context budget stayed trivial throughout (0.75%-8.28% of the 32,768-token window) -- this corpus is far too small to ever stress it.",
        col_widths=[1.8, 4.8, 2.6, 2.5],
    )

    # 16 -- real finding: missing document
    add_bullet_slide(
        prs, "Real Finding", "A Document Can Silently Vanish From the Answer",
        [
            (0, "synthetic_sample doesn't appear in the summary at all until top_k=10 -- more than half the entire corpus"),
            (0, "Not a bug: that document's table- and chart-heavy chunks are genuinely further in embedding space from a broad, prose-style question than the other two documents' chunks are"),
            (1, "A low top_k doesn't just retrieve less content -- it can drop an entire document from the answer with no error, no warning, nothing that looks wrong until you actually check which sources got cited."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 17 -- real finding: generation time explained
    add_bullet_slide(
        prs, "Real Finding", "51 Seconds Isn't a Glitch -- It's a Longer Answer",
        [
            (0, "51.4s at top_k=10 looks like an outlier next to 2-8s everywhere else, until you read the actual output"),
            (0, "At top_k=10, the model chose to write a longer, 5-part itemized answer -- long enough to hit the 250-token cap and get cut off mid-sentence"),
            (0, "Every other setting got a 1-3 sentence answer that stopped naturally well under the cap"),
            (1, "More retrieved context didn't just change WHAT got covered -- it changed the SHAPE of the answer the model chose to write."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 18 -- fun fact / honest finding: language leakage
    add_fact_slide(
        prs, "A Real Model Quirk", "Small Models Can Leak Stray Tokens Mid-Sentence",
        [
            (0, "In 2 of the 5 experiment runs (non-thinking mode), a stray non-English character or garbled word leaked into an otherwise-clean English answer"),
            (1, "\"...southern border d̯.\"   and   \"NASA ResearchGrammar\""),
            (0, "Real, repeatable, not cherry-picked -- worth knowing before trusting a 0.6B model's output unsupervised in anything user-facing"),
        ],
        next_page(),
    )

    # 19 -- real code: the prompt itself
    add_code_slide(
        prs, "The Actual Prompt", "Numbered, Source-Attributed Context",
        "rag/prompting.py::build_messages()",
        "[1] (source: nasa_iss_factsheet > Space Station Overview)\n"
        "It's also home to astronauts and cosmonauts...\n\n"
        "[2] (source: federal_register_proclamation > A Proclamation)\n"
        "...\n\n"
        "Question: What does this collection of documents cover?",
        "The Streamlit app shows this exact assembled text in an expander -- if an answer looks wrong, check what the model actually saw before assuming the model is broken.",
        next_page(),
        height_in=2.6,
    )

    # 20 -- the app
    add_bullet_slide(
        prs, "The Deliverable", "A Streamlit App, Not Just a Script",
        [
            (0, "Question box, real Retrieve + Summarize button"),
            (0, "Sidebar: top_k slider, embedding dimension selector, document filter, enable_thinking toggle, max_new_tokens slider, temperature slider (pre-filled with Qwen3's own recommended preset for the selected mode)"),
            (0, "Retrieved chunks shown with their real distance scores, before the summary"),
            (0, "The exact assembled prompt shown in an expander, plus prompt token count and % of the context window used"),
            (0, "Generation time, output token count, and the exact sampling parameters used, printed alongside every answer"),
        ],
        next_page(),
    )

    # 21 -- wrap-up
    add_bullet_slide(
        prs, "Wrap-Up", "What Module 6 Actually Completes",
        [
            (0, "A full pipeline, end to end: PDF -> classified regions -> validated markdown -> chunks -> embeddings in Postgres -> retrieval -> a real LLM answer, in a UI"),
            (0, "Two model-config details (enable_thinking, mode-specific sampling) implemented the way the model card actually specifies, not guessed"),
            (0, "A measured, not asserted, answer to \"does top_k matter\": yes, enough to silently drop an entire document from the answer"),
            (0, "Honest findings throughout, including the model's own real mistakes -- language leakage, an answer cut off mid-sentence -- not smoothed over"),
        ],
        next_page(),
    )

    # 22 -- closing
    add_closing_slide(
        prs,
        title="Questions?",
        subtitle="6. SUMMARISATION/  --  README.md has the full parameter reference and every real number in this deck.",
        footer="Module 6  ·  Data Processing → Analysis → Chunking → Embedding & Database → Database Deep Dive → Summarization",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
