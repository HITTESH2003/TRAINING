"""Builds Chunking_Module_Overview.pptx -- a teaching deck for Module 3 (Chunking).

Every number in this deck comes from actually running chunk_document.py
against Module 1's real output (3 documents, as of this writing) -- not
illustrative placeholders. Re-run chunk_document.py first if you want the
deck to reflect a fresh run; this script doesn't re-run it for you.

Same visual design system as 1. OCR/make_slides.py and
2. ANALYSIS/make_slides.py, duplicated here rather than imported -- each
module stays self-contained and only shares the venv.

Usage:
  python make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "Chunking_Module_Overview.pptx"

# ---- palette (matches the other two modules' decks) -------------------------
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x27)
MUTED = RGBColor(0x8A, 0x93, 0xA6)
MUTED_DARK = RGBColor(0x5B, 0x63, 0x74)
ACCENT = RGBColor(0x3E, 0x8E, 0xDE)       # blue -- structural/technical content
ACCENT_WARM = RGBColor(0xF2, 0xA6, 0x3D)  # amber -- "did you know" slides
ACCENT_GOOD = RGBColor(0x4C, 0xAF, 0x7D)  # green -- real code slides
ACCENT_BAD = RGBColor(0xE0, 0x6C, 0x5C)   # red -- real findings / case studies
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT = RGBColor(0xA6, 0xE3, 0xA1)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def _kicker(slide, text, color=ACCENT):
    _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5), text.upper(), 15, color, bold=True)


def _accent_bar(slide, color, top=Inches(0), height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _page_number(slide, n, dark=False):
    color = MUTED if dark else MUTED_DARK
    _textbox(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35), str(n), 11, color, align=PP_ALIGN.RIGHT)


# ---- adaptive bullet sizing (fits content instead of guessing a fixed size) --

_BULLET_SIZE_PRESETS = [
    (20, 16, 12, 6),
    (18, 15, 10, 5),
    (16, 13, 8, 4),
    (14.5, 12, 6, 3),
]


def _estimate_block_height_in(bullets, box_w_in, size0, size1, space0, space1, prefix0_len=3, prefix1_len=8):
    total = 0.0
    for level, text in bullets:
        size = size0 if level == 0 else size1
        space = space0 if level == 0 else space1
        prefix_len = prefix0_len if level == 0 else prefix1_len
        chars_per_line = max(10, int(box_w_in / (0.5 * size / 72)))
        lines = max(1, -(-(len(text) + prefix_len) // chars_per_line))
        line_height_in = (size * 1.22) / 72
        total += lines * line_height_in + space / 72
    return total


def _pick_bullet_sizes(bullets, box_w_in, box_h_in, presets=_BULLET_SIZE_PRESETS, safety_margin=0.88):
    # Require the estimate to clear a margin, not just technically fit --
    # the char-count-based line-wrap estimate is approximate, so an exact
    # fit on paper can still overflow once real font metrics render it.
    budget = box_h_in * safety_margin
    for preset in presets:
        if _estimate_block_height_in(bullets, box_w_in, *preset) <= budget:
            return preset
    return presets[-1]


def _fill_bullets(tf, bullets, size0, size1, space0, space1):
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space0 if level == 0 else space1)
        prefix = "▸  " if level == 0 else "     –  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size0 if level == 0 else size1)
        run.font.color.rgb = INK if level == 0 else MUTED_DARK
        run.font.name = FONT_BODY
        run.font.bold = level == 0


# ---- slide builders ---------------------------------------------------------

def add_title_slide(prs, kicker, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.5), kicker.upper(), 18, ACCENT, bold=True)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6), title, 44, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.35), Inches(11.3), Inches(1.0), subtitle, 20, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


def add_bullet_slide(prs, kicker, title, bullets, page, color=ACCENT, note=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.9), title, 30, INK, bold=True)

    box_w_in, box_h_in = 11.7, 4.9
    sizes = _pick_bullet_sizes(bullets, box_w_in, box_h_in)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(box_w_in), Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    _fill_bullets(tf, bullets, *sizes)

    if note:
        _textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45), note, 12, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_fact_slide(prs, kicker, title, bullets, page, color=ACCENT_WARM):
    slide = add_bullet_slide(prs, kicker, title, bullets, page, color=color)
    tag = slide.shapes.add_textbox(Inches(10.6), Inches(0.42), Inches(2.0), Inches(0.5))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "★ DID YOU KNOW"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_WARM
    return slide


def add_code_slide(prs, kicker, title, label, code_text, note, page, height_in=4.0):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT_GOOD)
    _kicker(slide, kicker, color=ACCENT_GOOD)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)
    _textbox(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4), label, 14, MUTED_DARK, bold=True)

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.7), Inches(height_in))
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line if line.strip() else " "
        run.font.size = Pt(14)
        run.font.name = FONT_CODE
        run.font.color.rgb = CODE_TEXT

    note_top = 2.0 + height_in + 0.2
    _textbox(slide, Inches(0.8), Inches(note_top), Inches(11.7), Inches(7.0 - note_top), note, 14, MUTED_DARK, italic=True)
    _page_number(slide, page)
    return slide


def add_table_slide(prs, kicker, title, headers, rows, page, color=ACCENT, note=None, col_widths=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.8), Inches(1.95), Inches(11.7), Inches(0.5 * n_rows)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table

    if col_widths:
        for c, w in enumerate(col_widths):
            table.columns[c].width = Inches(w)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = WHITE
                run.font.name = FONT_BODY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF2, 0xF4, 0xF8)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = INK
                    run.font.name = FONT_BODY

    if note:
        note_top = top + height + Inches(0.25)
        _textbox(slide, Inches(0.8), note_top, Inches(11.7), Inches(7.2) - note_top, note, 13, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_closing_slide(prs, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2), title, 40, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(1.4), subtitle, 18, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


SERIES_MODULES = [
    (1, "OCR & Document Understanding", "PDF to structured markdown"),
    (2, "Analysis & Validation", "Is the extraction actually right?"),
    (3, "Chunking", "Splitting documents for retrieval"),
    (4, "Embedding & Database", "Vectors, Postgres, pgvector"),
    (5, "Database Deep Dive", "Search, indexes, and what Postgres actually does"),
    (6, "Summarization", "Retrieval + a real LLM + a UI"),
]


def add_series_intro_slide(prs, current_module, page):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT)
    _kicker(slide, "Video Series")
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(1.1), "Part of a Series: How to Build Your First Chatbot", 28, INK, bold=True)

    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "This module is one step in a planned video series that builds a real, working "
        "chatbot end to end -- real documents, real models, real measured results, not toy examples."
    )
    run.font.size = Pt(16)
    run.font.color.rgb = MUTED_DARK
    run.font.name = FONT_BODY
    run.font.italic = True

    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(11.7), Inches(3.0))
    tf2 = list_box.text_frame
    tf2.word_wrap = True
    first = True
    for num, name, desc in SERIES_MODULES:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(10)
        is_current = num == current_module
        run = p.add_run()
        prefix = "►  " if is_current else "    "
        suffix = "   ◂ you are here" if is_current else ""
        run.text = f"{prefix}Module {num}: {name} -- {desc}{suffix}"
        run.font.size = Pt(18) if is_current else Pt(16)
        run.font.bold = is_current
        run.font.color.rgb = ACCENT if is_current else MUTED_DARK
        run.font.name = FONT_BODY

    _textbox(
        slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
        "More modules are planned as the series continues toward a complete chatbot -- retrieval and response generation are next.",
        13, MUTED_DARK, italic=True,
    )

    _page_number(slide, page)
    return slide


# ---- deck content -----------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = [0]

    def next_page():
        page[0] += 1
        return page[0]

    # 1 -- title
    add_title_slide(
        prs,
        kicker="Module 3 · Chunking",
        title="Splitting a Document for Retrieval",
        subtitle="Three genuinely different techniques, compared side by side on the same real documents Modules 1 and 2 already produced.",
        footer="recursive  ·  semantic (Qwen3-Embedding-0.6B)  ·  hierarchical",
    )
    next_page()

    # 2 -- series intro
    add_series_intro_slide(prs, current_module=3, page=next_page())

    # 3 -- why chunking
    add_bullet_slide(
        prs, "Why This Exists", "You Can't Just Hand Over the Whole Document",
        [
            (0, "An LLM prompt has a token budget. An embedding model has an input limit. Neither wants (or can take) an entire multi-page document at once"),
            (0, "So every retrieval system needs to cut documents into pieces first -- and where you cut matters as much as the content itself"),
            (0, "Cut in the wrong place and you get a chunk that's half a table, or a paragraph missing the sentence that explains it"),
            (0, "This module builds three different answers to “where do I cut?” and runs all three on the same documents so the differences are visible, not theoretical"),
        ],
        next_page(),
    )

    # 3 -- fun fact: RAG history
    add_fact_slide(
        prs, "Did you know", "“Chunking” Is Younger Than the Rest of This Pipeline's Ideas",
        [
            (0, "Retrieval-Augmented Generation (RAG) -- the technique that made “chunk it, embed it, retrieve it” a standard pattern -- comes from a 2020 Facebook AI Research paper (Lewis et al.)"),
            (0, "Before that, “split the document” mostly meant fixed-size windows for search indexing, not for feeding a generative model"),
            (0, "The field still hasn't converged on one right answer -- which is exactly why comparing strategies side by side, like this module does, is more honest than picking one and asserting it's best"),
        ],
        next_page(),
    )

    # 4 -- the core teaching point
    add_bullet_slide(
        prs, "The Core Idea", "Three Ideas, Not Four Flavors of One",
        [
            (0, "MECHANICS -- given a token size, chunk to it, respecting paragraph/sentence boundaries. No understanding of meaning or structure.  →  recursive"),
            (0, "MEANING -- group content by how similar it actually is, discovered from the content itself via embeddings. No reference to headings at all.  →  semantic"),
            (0, "STRUCTURE -- headings define what belongs under what; tables/figures are atomic; boundaries within a section follow how much content is really there.  →  hierarchical"),
            (0, "Each one answers “where should this chunk end?” with different evidence: a number, a similarity score, or a heading"),
        ],
        next_page(),
    )

    # 5 -- strategy 1: recursive
    add_bullet_slide(
        prs, "Strategy 1 of 3", "recursive — Mechanical, Token-Bounded",
        [
            (0, "Given a token size, that's what gets chunked -- the “just give me chunks of size N” baseline every RAG tutorial starts with"),
            (0, "Tries paragraph breaks first; if a piece is still too big, recurses into it with a finer separator: line break → sentence → word"),
            (0, "Same idea behind LangChain's RecursiveCharacterTextSplitter -- widely used precisely because it's simple and dependency-free"),
            (0, "Sizing is checked with the real tokenizer on every candidate piece, not an approximated character budget -- max_tokens is a real guarantee here"),
            (0, "Has zero idea what a heading, table, or topic is -- a long <table> gets cut exactly like a long paragraph would"),
        ],
        next_page(),
    )

    # 6 -- fun fact: recursive splitting adoption
    add_fact_slide(
        prs, "Did you know", "The “Naive” Baseline Is Still the Most Common Choice",
        [
            (0, "Despite every RAG blog post explaining why fixed/recursive splitting misses structure, it remains the default in most production retrieval pipelines"),
            (0, "Why: it's fast, deterministic, needs no extra model, and works “well enough” when documents are short and mostly plain text"),
            (0, "It breaks down exactly where this project's documents do -- tables, multi-section structure, mixed content types. Simple documents hide the problem; real ones expose it"),
        ],
        next_page(),
    )

    # 7 -- strategy 2: semantic
    add_bullet_slide(
        prs, "Strategy 2 of 3", "semantic — Meaning-Driven, Structure-Agnostic",
        [
            (0, "Embeds every block with Qwen3-Embedding-0.6B -- same model family as the rest of the pipeline, no new dependency, ~1.2GB download through transformers already installed"),
            (0, "Compares each block to the one before it: cosine similarity of L2-normalized vectors, which is just a dot product"),
            (0, "When similarity drops below a threshold (default 0.5), that's a topic shift -- a new chunk starts there, found from the content itself, with no “##” heading involved at all"),
            (0, "Every chunk records the exact boundary_similarity that ended it -- “why did it split here” has a real number behind it, not a guess"),
            (0, "Tables and figures still stay atomic -- they don't have a “topic” to be similar or dissimilar to"),
        ],
        next_page(),
    )

    # 8 -- real code: embedding_util.py
    add_code_slide(
        prs, "The Actual Code", "Cosine Similarity Is Just a Dot Product",
        "chunking/embedding_util.py",
        "# embeddings are L2-normalized on the way out, so:\n"
        "def cosine_similarity(a, b) -> float:\n"
        "    return float((a * b).sum())\n\n"
        "# a topic shift is detected as:\n"
        "similarity = cosine_similarity(prev_embedding, embedding)\n"
        "topic_shift = similarity < similarity_threshold   # default 0.5",
        "No separate similarity library needed -- once vectors are normalized, similarity is one line of arithmetic.",
        next_page(),
        height_in=3.2,
    )

    # 9 -- fun fact: embeddings history
    add_fact_slide(
        prs, "Did you know", "“king − man + woman ≈ queen” Is a Real, Famous Result",
        [
            (0, "word2vec (2013, Mikolov et al. at Google) was the paper that made this demo famous -- vector arithmetic on word embeddings actually landing on meaningful analogies"),
            (0, "That was single-word embeddings. Sentence/paragraph-level embeddings (what semantic.py uses) are the same underlying idea applied to whole chunks of meaning, not individual words"),
            (0, "Qwen3-Embedding-0.6B produces 1024-dimensional vectors -- every block in a document becomes a single point in a 1024-dimensional space, and “topic shift” just means two points are far apart"),
        ],
        next_page(),
    )

    # 10 -- strategy 3: hierarchical
    add_bullet_slide(
        prs, "Strategy 3 of 3", "hierarchical — Structure-Driven, Content-Sized",
        [
            (0, "Builds a section tree from real headings -- every chunk carries a section_path breadcrumb"),
            (0, "Tables and figures are atomic blocks, never split, whatever else is happening around them"),
            (0, "Packs paragraphs into each section's chunks up to a token budget -- an oversized single paragraph falls back to recursive.py's sentence-level splitter, and small trailing fragments get merged into their neighbor instead of standing alone"),
            (0, "Every chunk also carries a contextualized_text field: “[Section Name]\\n<chunk text>” -- prepending the section breadcrumb before embedding is a real retrieval technique (“contextual chunk headers”), not just described here but actually in the output"),
        ],
        next_page(),
    )

    # 11 -- connecting back to module 1
    add_bullet_slide(
        prs, "Connecting the Pipeline", "Real Layout Data, Not Re-Guessed Structure",
        [
            (0, "semantic and hierarchical both need to know what a heading, table, or figure IS before making structure-aware decisions"),
            (0, "An earlier version of this module regex-parsed the rendered markdown to re-guess that -- and silently lost information: table_caption, table_footnote, and formula_caption all render as plain or italic text, indistinguishable from a paragraph"),
            (0, "Module 1 now writes regions.json -- the real per-region layout record (true class, confidence score, page, bbox) -- and layout_blocks.py reads that directly instead"),
            (1, "This is exactly why Module 1 exists as a separate, inspectable step: its output became load-bearing for a module built weeks later"),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 12 -- real case study: table atomicity
    add_bullet_slide(
        prs, "Real Finding", "Two Strategies That Agree on Nothing Else, Agree on This",
        [
            (0, "On synthetic_sample.pdf: recursive.json splits the table's <table border=\"1\">... opening (198 tokens) from its closing </table> (70 tokens) -- both flagged table_split"),
            (0, "hierarchical.json AND semantic.json both keep the entire 251-token table as one oversized_atomic chunk instead -- over budget, but never broken"),
            (1, "One strategy decides boundaries from headings. The other decides them from embedding similarity. They don't agree on much -- except this, because atomicity is a rule applied BEFORE either strategy's actual logic runs"),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 13 -- real case study: low confidence flag
    add_bullet_slide(
        prs, "Real Finding", "A Confidence Score That Survives All the Way to a Chunk",
        [
            (0, "One region in synthetic_sample.pdf is a genuinely bad OCR result -- an overlapping/duplicate box the VLM transcribed as “q3 regional output” at confidence 0.30 (documented in Module 1's README)"),
            (0, "Any semantic or hierarchical chunk that includes that region is flagged contains_low_confidence_region"),
            (1, "That flag only exists because layout_blocks.py reads Module 1's real per-region confidence scores -- there's no way to reconstruct it from rendered markdown text alone"),
        ],
        next_page(),
    )

    # 14 -- real finding: semantic's honest limitation
    add_fact_slide(
        prs, "Honest Limitation", "Semantic Chunking Has No Small-Chunk Safety Net",
        [
            (0, "On nasa_iss_factsheet.pdf: semantic produces 8 chunks to hierarchical's 4 -- several under the 40-token minimum, including a genuine 3-token chunk (a duplicate running-header fragment)"),
            (0, "Every one of those splits has a real similarity score behind it -- a title block, an address block, and a repeated header really are three different “topics” to an embedding model"),
            (0, "hierarchical.py has an explicit merge step for small trailing fragments. semantic.py deliberately doesn't -- overriding a detected topic boundary just because the result is small would be a different, unstated rule sneaking back in"),
            (1, "A real, current design tradeoff, not an oversight -- worth debating either way"),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 15 -- real numbers table
    add_table_slide(
        prs, "Real Numbers", "Chunk Count & Size Spread, Same Documents",
        ["Document", "Strategy", "Chunks", "Avg Tokens", "Stdev"],
        [
            ["federal_register_proclamation", "recursive", 13, "116.5", "54.8"],
            ["federal_register_proclamation", "semantic", 13, "114.2", "45.4"],
            ["federal_register_proclamation", "hierarchical", 11, "131.1", "56.8"],
            ["nasa_iss_factsheet", "recursive", 4, "123.0", "71.1"],
            ["nasa_iss_factsheet", "semantic", 8, "58.0", "54.5"],
            ["nasa_iss_factsheet", "hierarchical", 4, "115.2", "68.0"],
        ],
        next_page(),
        note="Full table (all 3 documents × 3 strategies, plus table_split_count, oversized_atomic_count, and avg_boundary_similarity) is in chunk_comparison_report.xlsx.",
        col_widths=[4.6, 2.2, 1.6, 1.7, 1.6],
    )

    # 16 -- the excel report
    add_bullet_slide(
        prs, "The Deliverable", "One Report, Every Strategy, Side by Side",
        [
            (0, "chunk_comparison_report.xlsx -- Comparison sheet, one row per document × strategy"),
            (1, "chunk_count, avg/min/max/stdev tokens, tiny_chunks_below_min, table_split_count, oversized_atomic_count, avg_boundary_similarity"),
            (0, "Chunks Detail sheet -- every individual chunk: section path, pages spanned, block-type composition, boundary similarity, notes, text preview"),
            (0, "Plus output/<doc-name>/<strategy>.json -- the actual chunks, ready to feed to an embedding step or an LLM prompt"),
        ],
        next_page(),
    )

    # 17 -- wrap-up
    add_bullet_slide(
        prs, "Wrap-Up", "What Module 3 Actually Produces",
        [
            (0, "Three real, runnable chunking strategies -- not descriptions of them, working code compared on the same documents"),
            (0, "A quantified argument for structure- and meaning-aware chunking over blind splitting -- table atomicity, confidence flagging, section breadcrumbs, all visible in real output"),
            (0, "An honest accounting of tradeoffs too -- semantic's fragmentation, hierarchical's reliance on headings existing at all"),
            (0, "This completes the pipeline: Module 1 extracts and classifies, Module 2 validates what came out, Module 3 turns it into retrieval-ready pieces"),
        ],
        next_page(),
    )

    # 18 -- closing
    add_closing_slide(
        prs,
        title="Questions?",
        subtitle="3. CHUNKING/  --  README.md has the full strategy reference and every real finding in this deck.",
        footer="Module 3 of 3  ·  Data Processing → Analysis → Chunking",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
