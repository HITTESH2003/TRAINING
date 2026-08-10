"""Builds Embedding_Database_Module_Overview.pptx -- a teaching deck for
Module 4 (Embedding & Database).

Every number in this deck comes from actually running load_and_store.py and
run_experiment.py against a real Postgres/pgvector container -- not
illustrative placeholders. Re-run both first if you want the deck to
reflect a fresh run; this script doesn't re-run them for you.

Same visual design system as the other three modules' decks, duplicated
here rather than imported -- each module stays self-contained and only
shares the venv.

Usage:
  python make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "Embedding_Database_Module_Overview.pptx"

# ---- palette (matches the other three modules' decks) -----------------------
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x27)
MUTED = RGBColor(0x8A, 0x93, 0xA6)
MUTED_DARK = RGBColor(0x5B, 0x63, 0x74)
ACCENT = RGBColor(0x3E, 0x8E, 0xDE)       # blue -- structural/technical content
ACCENT_WARM = RGBColor(0xF2, 0xA6, 0x3D)  # amber -- "did you know" slides
ACCENT_GOOD = RGBColor(0x4C, 0xAF, 0x7D)  # green -- real code slides
ACCENT_BAD = RGBColor(0xE0, 0x6C, 0x5C)   # red -- real findings / case studies
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT = RGBColor(0xA6, 0xE3, 0xA1)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def _kicker(slide, text, color=ACCENT):
    _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5), text.upper(), 15, color, bold=True)


def _accent_bar(slide, color, top=Inches(0), height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _page_number(slide, n, dark=False):
    color = MUTED if dark else MUTED_DARK
    _textbox(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35), str(n), 11, color, align=PP_ALIGN.RIGHT)


# ---- adaptive bullet sizing (fits content instead of guessing a fixed size) --

_BULLET_SIZE_PRESETS = [
    (20, 16, 12, 6),
    (18, 15, 10, 5),
    (16, 13, 8, 4),
    (14.5, 12, 6, 3),
]


def _estimate_block_height_in(bullets, box_w_in, size0, size1, space0, space1, prefix0_len=3, prefix1_len=8):
    total = 0.0
    for level, text in bullets:
        size = size0 if level == 0 else size1
        space = space0 if level == 0 else space1
        prefix_len = prefix0_len if level == 0 else prefix1_len
        chars_per_line = max(10, int(box_w_in / (0.5 * size / 72)))
        lines = max(1, -(-(len(text) + prefix_len) // chars_per_line))
        line_height_in = (size * 1.22) / 72
        total += lines * line_height_in + space / 72
    return total


def _pick_bullet_sizes(bullets, box_w_in, box_h_in, presets=_BULLET_SIZE_PRESETS, safety_margin=0.88):
    # Require the estimate to clear a margin, not just technically fit --
    # the char-count-based line-wrap estimate is approximate, so an exact
    # fit on paper can still overflow once real font metrics render it.
    budget = box_h_in * safety_margin
    for preset in presets:
        if _estimate_block_height_in(bullets, box_w_in, *preset) <= budget:
            return preset
    return presets[-1]


def _fill_bullets(tf, bullets, size0, size1, space0, space1):
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space0 if level == 0 else space1)
        prefix = "▸  " if level == 0 else "     –  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size0 if level == 0 else size1)
        run.font.color.rgb = INK if level == 0 else MUTED_DARK
        run.font.name = FONT_BODY
        run.font.bold = level == 0


# ---- slide builders ---------------------------------------------------------

def add_title_slide(prs, kicker, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.5), kicker.upper(), 18, ACCENT, bold=True)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6), title, 44, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.35), Inches(11.3), Inches(1.0), subtitle, 20, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


def add_bullet_slide(prs, kicker, title, bullets, page, color=ACCENT, note=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.9), title, 30, INK, bold=True)

    box_w_in, box_h_in = 11.7, 4.9
    sizes = _pick_bullet_sizes(bullets, box_w_in, box_h_in)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(box_w_in), Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    _fill_bullets(tf, bullets, *sizes)

    if note:
        _textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45), note, 12, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_fact_slide(prs, kicker, title, bullets, page, color=ACCENT_WARM):
    slide = add_bullet_slide(prs, kicker, title, bullets, page, color=color)
    tag = slide.shapes.add_textbox(Inches(10.6), Inches(0.42), Inches(2.0), Inches(0.5))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "★ DID YOU KNOW"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_WARM
    return slide


def add_code_slide(prs, kicker, title, label, code_text, note, page, height_in=4.0):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT_GOOD)
    _kicker(slide, kicker, color=ACCENT_GOOD)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)
    _textbox(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4), label, 14, MUTED_DARK, bold=True)

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.7), Inches(height_in))
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line if line.strip() else " "
        run.font.size = Pt(14)
        run.font.name = FONT_CODE
        run.font.color.rgb = CODE_TEXT

    note_top = 2.0 + height_in + 0.2
    _textbox(slide, Inches(0.8), Inches(note_top), Inches(11.7), Inches(7.0 - note_top), note, 14, MUTED_DARK, italic=True)
    _page_number(slide, page)
    return slide


def add_table_slide(prs, kicker, title, headers, rows, page, color=ACCENT, note=None, col_widths=None):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.8), Inches(1.95), Inches(11.7), Inches(0.5 * n_rows)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table

    if col_widths:
        for c, w in enumerate(col_widths):
            table.columns[c].width = Inches(w)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = WHITE
                run.font.name = FONT_BODY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF2, 0xF4, 0xF8)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = INK
                    run.font.name = FONT_BODY

    if note:
        note_top = top + height + Inches(0.25)
        _textbox(slide, Inches(0.8), note_top, Inches(11.7), Inches(7.2) - note_top, note, 13, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_closing_slide(prs, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2), title, 40, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(1.4), subtitle, 18, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


SERIES_MODULES = [
    (1, "OCR & Document Understanding", "PDF to structured markdown"),
    (2, "Analysis & Validation", "Is the extraction actually right?"),
    (3, "Chunking", "Splitting documents for retrieval"),
    (4, "Embedding & Database", "Vectors, Postgres, pgvector"),
    (5, "Database Deep Dive", "Search, indexes, and what Postgres actually does"),
    (6, "Summarization", "Retrieval + a real LLM + a UI"),
]


def add_series_intro_slide(prs, current_module, page):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT)
    _kicker(slide, "Video Series")
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(1.1), "Part of a Series: How to Build Your First Chatbot", 28, INK, bold=True)

    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "This module is one step in a planned video series that builds a real, working "
        "chatbot end to end -- real documents, real models, real measured results, not toy examples."
    )
    run.font.size = Pt(16)
    run.font.color.rgb = MUTED_DARK
    run.font.name = FONT_BODY
    run.font.italic = True

    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(11.7), Inches(3.0))
    tf2 = list_box.text_frame
    tf2.word_wrap = True
    first = True
    for num, name, desc in SERIES_MODULES:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(10)
        is_current = num == current_module
        run = p.add_run()
        prefix = "►  " if is_current else "    "
        suffix = "   ◂ you are here" if is_current else ""
        run.text = f"{prefix}Module {num}: {name} -- {desc}{suffix}"
        run.font.size = Pt(18) if is_current else Pt(16)
        run.font.bold = is_current
        run.font.color.rgb = ACCENT if is_current else MUTED_DARK
        run.font.name = FONT_BODY

    _textbox(
        slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
        "More modules are planned as the series continues toward a complete chatbot -- retrieval and response generation are next.",
        13, MUTED_DARK, italic=True,
    )

    _page_number(slide, page)
    return slide


# ---- deck content -----------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = [0]

    def next_page():
        page[0] += 1
        return page[0]

    # 1 -- title
    add_title_slide(
        prs,
        kicker="Module 4 · Embedding & Database",
        title="Embedding & Storing at Scale",
        subtitle="Module 3's chunks become searchable vectors -- nomic-embed-text-v1.5, Postgres + pgvector, and a real measured answer to 'does dimension size actually matter?'",
        footer="nomic-embed-text-v1.5  ·  pgvector/pgvector:pg16  ·  Docker",
    )
    next_page()

    # 2 -- series intro
    add_series_intro_slide(prs, current_module=4, page=next_page())

    # 3 -- why this module
    add_bullet_slide(
        prs, "Why This Exists", "From Chunks to a Queryable System",
        [
            (0, "Module 3 produced chunks -- structured pieces of text. A chunk sitting in a JSON file isn't searchable by meaning yet"),
            (0, "To search by meaning, every chunk needs to become a vector (an embedding), and those vectors need somewhere real to live -- not just in memory, in a database built for this"),
            (0, "This module does both, and then asks a question most tutorials skip: the embedding model can produce smaller vectors -- does that actually cost you accuracy, or is it free?"),
            (1, "Docker is only here as \"the way we run Postgres.\" This isn't a Docker module -- that's coming later."),
        ],
        next_page(),
    )

    # 3 -- fun fact: vector search is old
    add_fact_slide(
        prs, "Did you know", "“Vector Search” Predates Every AI Buzzword Attached to It",
        [
            (0, "The core idea -- represent a document as a vector, rank others by how close their vectors are -- is the Vector Space Model, published by Gerard Salton in 1975, for classic information retrieval, decades before embeddings or neural networks"),
            (0, "pgvector itself (the Postgres extension this module uses) was released in 2021 by a single developer, Andrew Kane -- not a big-lab research project"),
            (0, "“Vector database” as a product category is newer still -- most of what it does is this 50-year-old idea, made fast enough to run at scale"),
        ],
        next_page(),
    )

    # 4 -- matryoshka concept
    add_bullet_slide(
        prs, "The Core Idea", "One Embedding, Many Valid Sizes",
        [
            (0, "nomic-embed-text-v1.5 is trained with Matryoshka Representation Learning: its native embedding is 768 numbers, but the FIRST 256 of those 768 numbers are themselves a valid, meaningful embedding on their own"),
            (0, "Same for the first 128, the first 64 -- nested representations, like the nesting dolls the technique is named for"),
            (0, "That's exactly what this module needs to test “dimension vs. accuracy” honestly: one model, one embedding per chunk, truncated to different sizes -- not five unrelated models that wouldn't be a fair comparison at all"),
        ],
        next_page(),
    )

    # 5 -- real code: truncation
    add_code_slide(
        prs, "The Actual Code", "Truncation Is Not “Just Slice the Vector”",
        "embedding/nomic_embed.py -- copied directly from the model's own card",
        "embeddings = F.layer_norm(embeddings, normalized_shape=(embeddings.shape[-1],))\n"
        "embeddings = embeddings[..., :dim]\n"
        "embeddings = F.normalize(embeddings, p=2, dim=-1)",
        "layer_norm, THEN slice, THEN re-normalize -- in that order. Slicing an already-L2-normalized 768-dim vector and renormalizing is a different operation the model was never trained for.",
        next_page(),
        height_in=2.2,
    )

    # 6 -- fun fact: prefixes
    add_fact_slide(
        prs, "Did you know", "This Model Fails Silently If You Skip One String",
        [
            (0, "Nomic trained this model to expect a task prefix on every input: \"search_document: \" on everything you store, \"search_query: \" on everything you search with"),
            (0, "Get it backwards, or skip it entirely, and nothing errors -- you get embeddings back, they look completely normal, and they're just measurably worse at retrieval"),
            (0, "This is a common, invisible bug in real RAG systems: a model \"working\" is not the same as a model being used correctly"),
        ],
        next_page(),
    )

    # 7 -- the database schema
    add_code_slide(
        prs, "The Storage", "One Table, One Column Per Dimension",
        "embedding/db.py -- actual schema, one row per chunk",
        "CREATE TABLE chunks (\n"
        "    chunk_id TEXT PRIMARY KEY,\n"
        "    doc_name TEXT NOT NULL,\n"
        "    section_path TEXT,\n"
        "    chunk_text TEXT NOT NULL,\n"
        "    embedding_768 vector(768),\n"
        "    embedding_512 vector(512),\n"
        "    embedding_256 vector(256),\n"
        "    embedding_128 vector(128),\n"
        "    embedding_64  vector(64)\n"
        ");",
        "Parallel columns, not five separate tables -- \"same chunk, different dimension\" is a column lookup, not a join.",
        next_page(),
        height_in=3.3,
    )

    # 8 -- real code: pgvector query
    add_code_slide(
        prs, "The Actual Search", "A Real SQL Query, Not Numpy",
        "embedding/db.py :: nearest_neighbors()",
        "SELECT chunk_id, doc_name, chunk_text,\n"
        "       embedding_256 <=> %s AS distance\n"
        "FROM chunks\n"
        "ORDER BY embedding_256 <=> %s\n"
        "LIMIT 5;",
        "<=> is pgvector's cosine-distance operator. 0 = identical direction, larger = less similar. This runs inside Postgres itself, not in Python.",
        next_page(),
        height_in=2.4,
    )

    # 9 -- fun fact: docker (kept light)
    add_fact_slide(
        prs, "Did you know", "Docker Started as an Internal Tool Nobody Planned to Open-Source",
        [
            (0, "Docker began in 2013 as an internal project at dotCloud, a struggling platform-as-a-service company -- it was open-sourced almost as an afterthought, to try to save the company"),
            (0, "The company pivoted entirely around Docker within a year and renamed itself Docker, Inc."),
            (0, "For this module, Docker's job is one line: give you a real Postgres without installing Postgres. That's the whole relationship -- the deeper Docker material is a later class."),
        ],
        next_page(),
    )

    # 10 -- experiment design
    add_bullet_slide(
        prs, "The Experiment", "18 Chunks, 10 Hand-Verified Questions",
        [
            (0, "Corpus: all 18 hierarchical chunks from Module 3, across all 3 documents -- the real chunks Module 3 actually produced, not a curated sample"),
            (0, "10 queries, each targeting one fact that appears in exactly one specific chunk -- every expected_chunk_id was picked by a human reading the real chunk text, verifiable by eye"),
            (0, "Each query runs through a real pgvector search at every one of the 5 test dimensions (768, 512, 256, 128, 64), and we record exactly where the correct chunk landed in the ranking"),
        ],
        next_page(),
    )

    # 11 -- real numbers table
    add_table_slide(
        prs, "Real Numbers", "Accuracy at Every Dimension",
        ["Dimension", "Recall@1", "Recall@3", "MRR", "Avg Distance to Correct"],
        [
            [768, "0.70", "1.00", "0.82", "0.2645"],
            [512, "0.70", "0.90", "0.79", "0.2671"],
            [256, "0.70", "0.90", "0.77", "0.2527"],
            [128, "0.70", "0.90", "0.77", "0.2253"],
            [64, "0.70", "0.90", "0.77", "0.1915"],
        ],
        next_page(),
        note="Recall@k = fraction of the 10 queries where the correct chunk landed in the top k results. Full detail (every query, every dimension) is in dimension_accuracy_report.xlsx.",
        col_widths=[2.0, 2.2, 2.2, 2.0, 3.3],
    )

    # 12 -- real finding: the degrading query
    add_bullet_slide(
        prs, "Real Finding", "The Aggregate Number Hides a Real Failure",
        [
            (0, "Recall@1 looks perfectly flat at 0.70 across every single dimension -- that's the number most people would report and move on"),
            (0, "One specific query tells a completely different story: “Who signed the regional operations report as Regional Director?”"),
            (1, "768 dims: rank 2   →   512 dims: rank 5   →   256 / 128 / 64 dims: not in the top 5 at all"),
            (0, "Every other query held roughly steady. This one -- whose answer is a short, low-information fact (a name and a title) -- degrades hard and monotonically as dimension shrinks"),
            (1, "The aggregate metric never shows this. Only the per-query detail does."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 13 -- real finding: the distance trap
    add_fact_slide(
        prs, "A Trap Worth Naming", "Smaller Distance Numbers Can Mean Worse Retrieval",
        [
            (0, "avg_distance_to_correct actually IMPROVES (gets smaller) as dimension shrinks: 0.2645 at 768 dims down to 0.1915 at 64 dims"),
            (0, "That looks like the model is “more confident” at lower dimensions. It isn't -- a coarser embedding space compresses everything closer together, correct answers and wrong ones alike"),
            (0, "The ranking metrics (recall@k, MRR) are what actually tell you whether retrieval held up. The raw similarity score can improve while the thing you actually care about gets worse"),
            (1, "Don't tune a real system by watching similarity scores go up -- watch whether the right answer keeps winning."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 14 -- connecting the pipeline
    add_bullet_slide(
        prs, "Connecting the Pipeline", "Every Module Feeds the Next One",
        [
            (0, "Module 1 turned a PDF into classified regions and clean markdown"),
            (0, "Module 2 turned that output into a validation report -- tokens, confidence, artifacts"),
            (0, "Module 3 turned it into structure-aware, retrieval-sized chunks"),
            (0, "Module 4 turned those chunks into a real, queryable, measured system -- and the ground truth for its own experiment came from actually reading Module 3's real chunk text, not a synthetic dataset"),
        ],
        next_page(),
    )

    # 15 -- wrap-up
    add_bullet_slide(
        prs, "Wrap-Up", "What Module 4 Actually Produces",
        [
            (0, "A real Postgres + pgvector database, 18 chunks, 5 embedding dimensions each, running in Docker"),
            (0, "A correct, model-card-verified embedding + truncation pipeline -- not a guess at how Matryoshka embeddings work"),
            (0, "A measured answer to “does dimension size matter”: mostly no, in aggregate -- but real, specific, predictable exceptions exist, and only per-query analysis finds them"),
            (0, "dimension_accuracy_report.xlsx -- every query, every dimension, every rank, ready to inspect"),
        ],
        next_page(),
    )

    # 16 -- closing
    add_closing_slide(
        prs,
        title="Questions?",
        subtitle="4. EMBEDDING AND DATABASE/  --  README.md has the full setup, the exact truncation code, and every real number in this deck.",
        footer="Module 4  ·  Data Processing → Analysis → Chunking → Embedding & Database",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
