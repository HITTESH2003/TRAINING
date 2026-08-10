"""Builds OCR_Module_Overview.pptx -- a teaching deck for Module 1 (OCR).

Every number, prompt, and finding in this deck is pulled from this project's
own README and real pipeline runs, not generic OCR trivia dressed up as
results. Run it after generating at least one real output (annotated
screenshots get embedded from output/nasa_iss_factsheet/ and
output/federal_register_proclamation/ if present; the deck still builds
without them, just without those two image slides).

Usage:
  python make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "OCR_Module_Overview.pptx"

# ---- palette --------------------------------------------------------------
DARK = RGBColor(0x0F, 0x17, 0x2A)
DARK_2 = RGBColor(0x16, 0x21, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x27)
MUTED = RGBColor(0x8A, 0x93, 0xA6)
MUTED_DARK = RGBColor(0x5B, 0x63, 0x74)
ACCENT = RGBColor(0x3E, 0x8E, 0xDE)      # blue -- structural/technical content
ACCENT_WARM = RGBColor(0xF2, 0xA6, 0x3D)  # amber -- "did you know" / fun-fact slides
ACCENT_GOOD = RGBColor(0x4C, 0xAF, 0x7D)  # green -- code/prompt slides
ACCENT_BAD = RGBColor(0xE0, 0x6C, 0x5C)   # red -- real-bug case studies
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT = RGBColor(0xA6, 0xE3, 0xA1)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def _kicker(slide, text, color=ACCENT, dark=False):
    _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5), text.upper(), 15, color, bold=True)


def _accent_bar(slide, color, top=Inches(0), height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _page_number(slide, n, dark=False):
    color = MUTED if dark else MUTED_DARK
    _textbox(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35), str(n), 11, color, align=PP_ALIGN.RIGHT)


# ---- slide builders ---------------------------------------------------------

def add_title_slide(prs, kicker, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.5), kicker.upper(), 18, ACCENT, bold=True)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6), title, 44, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.35), Inches(11.3), Inches(1.0), subtitle, 20, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


def add_section_slide(prs, n, number, title, subtitle, page):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(0.9), Inches(2.4), Inches(3.5), Inches(2), number, 96, ACCENT, bold=True)
    _textbox(slide, Inches(1), Inches(4.3), Inches(11), Inches(1), title, 36, WHITE, bold=True)
    if subtitle:
        _textbox(slide, Inches(1), Inches(5.15), Inches(10.5), Inches(1), subtitle, 18, MUTED)
    _page_number(slide, page, dark=True)
    return slide


# Progressively smaller (size0, size1, space0, space1) presets to try, largest first --
# lets a slide with few short bullets render big and confident, and a slide with a lot
# to say shrink just enough to fit rather than overflow the box.
_BULLET_SIZE_PRESETS = [
    (20, 16, 12, 6),
    (18, 15, 10, 5),
    (16, 13, 8, 4),
    (14.5, 12, 6, 3),
]


def _estimate_block_height_in(bullets, box_w_in, size0, size1, space0, space1, prefix0_len=3, prefix1_len=8):
    total = 0.0
    for level, text in bullets:
        size = size0 if level == 0 else size1
        space = space0 if level == 0 else space1
        prefix_len = prefix0_len if level == 0 else prefix1_len
        chars_per_line = max(10, int(box_w_in / (0.5 * size / 72)))
        lines = max(1, -(-(len(text) + prefix_len) // chars_per_line))
        line_height_in = (size * 1.22) / 72
        total += lines * line_height_in + space / 72
    return total


def _pick_bullet_sizes(bullets, box_w_in, box_h_in, safety_margin=0.88):
    # Require the estimate to clear a margin, not just technically fit --
    # the char-count-based line-wrap estimate is approximate, so an exact
    # fit on paper can still overflow once real font metrics render it.
    budget = box_h_in * safety_margin
    for preset in _BULLET_SIZE_PRESETS:
        if _estimate_block_height_in(bullets, box_w_in, *preset) <= budget:
            return preset
    return _BULLET_SIZE_PRESETS[-1]


def _fill_bullets(tf, bullets, size0, size1, space0, space1):
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space0 if level == 0 else space1)
        prefix = "▸  " if level == 0 else "     –  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size0 if level == 0 else size1)
        run.font.color.rgb = INK if level == 0 else MUTED_DARK
        run.font.name = FONT_BODY
        run.font.bold = level == 0


def add_bullet_slide(prs, kicker, title, bullets, page, color=ACCENT, note=None):
    """bullets: list of (level, text) tuples, level 0 or 1."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.9), title, 30, INK, bold=True)

    box_w_in, box_h_in = 11.7, 4.9
    sizes = _pick_bullet_sizes(bullets, box_w_in, box_h_in)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(box_w_in), Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    _fill_bullets(tf, bullets, *sizes)

    if note:
        _textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45), note, 12, MUTED_DARK, italic=True)

    _page_number(slide, page)
    return slide


def add_fact_slide(prs, kicker, title, bullets, page, color=ACCENT_WARM):
    slide = add_bullet_slide(prs, kicker, title, bullets, page, color=color)
    # amber corner tag to make "trivia" slides visually distinct at a glance
    tag = slide.shapes.add_textbox(Inches(10.6), Inches(0.42), Inches(2.0), Inches(0.5))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "★ DID YOU KNOW"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_WARM
    return slide


def add_prompt_slide(prs, kicker, title, label, prompt_text, note, page):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT_GOOD)
    _kicker(slide, kicker, color=ACCENT_GOOD)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)
    _textbox(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.4), label, 14, MUTED_DARK, bold=True)

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.9))
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = prompt_text
    run.font.size = Pt(15)
    run.font.name = FONT_CODE
    run.font.color.rgb = CODE_TEXT

    _textbox(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.0), note, 14, MUTED_DARK, italic=True)
    _page_number(slide, page)
    return slide


def add_image_slide(prs, kicker, title, image_path, bullets, page, color=ACCENT):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, color)
    _kicker(slide, kicker, color=color)
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7), title, 28, INK, bold=True)

    # image on the left, notes on the right
    if image_path and image_path.exists():
        from PIL import Image as PILImage

        with PILImage.open(image_path) as im:
            iw, ih = im.size
        max_w, max_h = Inches(6.6), Inches(5.0)
        scale = min(max_w / iw, max_h / ih)
        w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
        left = Inches(0.7)
        top = Inches(1.75) + (max_h - h) / 2
        slide.shapes.add_picture(str(image_path), left, top, width=w, height=h)
        text_left = Inches(7.6)
        text_width = Inches(5.0)
    else:
        text_left = Inches(0.8)
        text_width = Inches(11.7)

    text_w_in, text_h_in = text_width / 914400, 4.8
    presets = [(16, 14, 14, 8), (14, 12.5, 10, 6), (12.5, 11, 8, 4)]
    sizes = None
    for preset in presets:
        if _estimate_block_height_in(bullets, text_w_in, *preset, prefix0_len=3, prefix1_len=6) <= text_h_in:
            sizes = preset
            break
    sizes = sizes or presets[-1]

    box = slide.shapes.add_textbox(text_left, Inches(1.9), text_width, Inches(text_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    _fill_bullets(tf, bullets, *sizes)

    _page_number(slide, page)
    return slide


def add_closing_slide(prs, title, subtitle, footer):
    slide = _blank(prs)
    _set_bg(slide, DARK)
    _accent_bar(slide, ACCENT)
    _textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2), title, 40, WHITE, bold=True)
    _textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(1.4), subtitle, 18, MUTED)
    _textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), footer, 13, MUTED_DARK)
    return slide


SERIES_MODULES = [
    (1, "OCR & Document Understanding", "PDF to structured markdown"),
    (2, "Analysis & Validation", "Is the extraction actually right?"),
    (3, "Chunking", "Splitting documents for retrieval"),
    (4, "Embedding & Database", "Vectors, Postgres, pgvector"),
    (5, "Database Deep Dive", "Search, indexes, and what Postgres actually does"),
    (6, "Summarization", "Retrieval + a real LLM + a UI"),
]


def add_series_intro_slide(prs, current_module, page):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _accent_bar(slide, ACCENT)
    _kicker(slide, "Video Series")
    _textbox(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(1.1), "Part of a Series: How to Build Your First Chatbot", 28, INK, bold=True)

    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "This module is one step in a planned video series that builds a real, working "
        "chatbot end to end -- real documents, real models, real measured results, not toy examples."
    )
    run.font.size = Pt(16)
    run.font.color.rgb = MUTED_DARK
    run.font.name = FONT_BODY
    run.font.italic = True

    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(11.7), Inches(3.0))
    tf2 = list_box.text_frame
    tf2.word_wrap = True
    first = True
    for num, name, desc in SERIES_MODULES:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(10)
        is_current = num == current_module
        run = p.add_run()
        prefix = "►  " if is_current else "    "
        suffix = "   ◂ you are here" if is_current else ""
        run.text = f"{prefix}Module {num}: {name} -- {desc}{suffix}"
        run.font.size = Pt(18) if is_current else Pt(16)
        run.font.bold = is_current
        run.font.color.rgb = ACCENT if is_current else MUTED_DARK
        run.font.name = FONT_BODY

    _textbox(
        slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
        "More modules are planned as the series continues toward a complete chatbot -- retrieval and response generation are next.",
        13, MUTED_DARK, italic=True,
    )

    _page_number(slide, page)
    return slide


# ---- deck content -----------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = [0]

    def next_page():
        page[0] += 1
        return page[0]

    annotated_nasa = ROOT / "output" / "nasa_iss_factsheet" / "annotated" / "page_001.png"
    annotated_fedreg = ROOT / "output" / "federal_register_proclamation" / "annotated" / "page_001.png"

    # 1 -- title
    add_title_slide(
        prs,
        kicker="Module 1 · Data Processing",
        title="OCR & Document Understanding",
        subtitle="From a raw PDF to structured markdown -- a small doc-layout model + a small vision-language model, working together.",
        footer="Qwen/Qwen3.5-0.8B  ·  DocLayout-YOLO  ·  PyMuPDF",
    )
    next_page()

    # 2 -- series intro
    add_series_intro_slide(prs, current_module=1, page=next_page())

    # 3 -- what is OCR, the old way vs the new way
    add_bullet_slide(
        prs, "Intro", "OCR Used to Mean One Thing",
        [
            (0, "Classic OCR: recognize characters. Feed it a scanned page, get back a string of text."),
            (1, "Built for clean, single-column, printed pages -- and it shows the moment a table or a two-column layout shows up"),
            (0, "It never asked “what kind of content is this?” -- a table and a paragraph were both just “characters on a page”"),
            (0, "Modern approach: separate “where is it and what is it” (layout detection) from “what does it say” (a vision-language model)"),
            (0, "That split is this entire module -- Steps 2-4 are exactly that division of labor"),
        ],
        next_page(),
    )

    # 3 -- fun facts: history of OCR
    add_fact_slide(
        prs, "Did you know", "OCR Has Been Around Longer Than You'd Think",
        [
            (0, "1929 -- Gustav Tauschek patents the first OCR device in Germany, using physical templates to match letterforms"),
            (0, "1974 -- Ray Kurzweil builds the first “omni-font” OCR reading machine for the blind; Stevie Wonder buys the first commercial unit"),
            (0, "Tesseract, the open-source OCR engine, was originally built at HP in the 1980s -- then shelved for a decade before Google open-sourced it in 2005"),
            (0, "For 90+ years, “OCR” meant character recognition. Vision-language models are the first real shift in what the term even means"),
        ],
        next_page(),
    )

    # 4 -- the real problem
    add_bullet_slide(
        prs, "The Problem", "A PDF Page Is Not Just Text",
        [
            (0, "A single page can contain: paragraphs, headings, tables, charts, photos -- and artifacts: logos, signatures, stamps/seals"),
            (0, "Naive text extraction flattens all of it into one stream -- a table becomes a wall of misaligned numbers, images just vanish"),
            (0, "This module's job: classify every region on the page BEFORE extracting anything from it"),
            (1, "Only then does it make sense to ask a table for its structure and a paragraph for its words, differently"),
        ],
        next_page(),
    )

    # 5 -- pipeline overview (section divider style but with the 7 steps listed)
    add_bullet_slide(
        prs, "Architecture", "Seven Steps, PDF to Markdown",
        [
            (0, "1. PDF → Image  --  render every page to a PNG"),
            (0, "2. Layout Detection  --  classify every region on the page"),
            (0, "3. Visualize  --  draw the boxes before doing anything else with them"),
            (0, "4. VLM Extraction  --  OCR / table-to-HTML / figure description, per region"),
            (0, "5. Region Routing  --  decide what belongs in the document body vs. metadata"),
            (0, "6. Assembly  --  reading order → one markdown file"),
            (0, "7. Post-processing + Metadata  --  clean up, validate, and record what happened"),
        ],
        next_page(),
    )

    # 6 -- step 1
    add_bullet_slide(
        prs, "Step 1 of 7", "PDF → Image",
        [
            (0, "PyMuPDF renders every page to a PNG at 200 DPI"),
            (0, "Why render to an image at all? Because the layout detector is a vision model -- it sees pixels, not the PDF's internal text objects"),
            (1, "This also means the pipeline works identically on a scanned photo of a document, not just a “digitally native” PDF"),
            (0, "Native PDF metadata (title, author, creation date, page count) is captured here too, straight from the file"),
        ],
        next_page(),
    )

    # 7 -- step 2 layout detection
    add_bullet_slide(
        prs, "Step 2 of 7", "Layout Detection — DocLayout-YOLO",
        [
            (0, "A YOLOv10-based detector, trained on DocStructBench (built on DocLayNet)"),
            (0, "Classifies every region into one of 10 classes:"),
            (1, "title, plain text, table, figure"),
            (1, "figure_caption, table_caption, table_footnote"),
            (1, "isolate_formula, formula_caption"),
            (1, "abandon -- headers, footers, watermarks"),
            (0, "Real run: 19 regions on page 1 of a 3-page Federal Register document, 15 on a NASA fact sheet"),
        ],
        next_page(),
    )

    # 8 -- image slide: bounding boxes on a real document
    add_image_slide(
        prs, "Step 3 of 7", "Seeing What the Model Actually Found",
        annotated_nasa,
        [
            (0, "Every box + label is drawn BEFORE any VLM call happens -- this is the layout detector's raw output, unedited"),
            (0, "Real finding: the NASA logo (top right) is boxed as “abandon” (gray), not “figure”"),
            (1, "The detector is trained mostly on academic papers -- a flat corporate logo doesn't look like anything in that training set"),
            (1, "It still gets logged -- just to metadata, not silently dropped"),
        ],
        next_page(),
    )

    # 9 -- step 4 VLM
    add_bullet_slide(
        prs, "Step 4 of 7", "The Vision-Language Model — Qwen3.5-0.8B",
        [
            (0, "0.8 billion parameters -- small enough to run on a laptop GPU (Apple MPS / CUDA), not a cloud API"),
            (0, "One model does three jobs, just with different prompts:"),
            (1, "Transcribe text (title, paragraphs, captions, formulas)"),
            (1, "Convert a table image into HTML"),
            (1, "Classify + describe a figure (chart, photo, logo, signature, stamp/seal...)"),
            (0, "The teaching tradeoff: a 0.8B model is small enough to inspect and reason about, and small enough to get things wrong in instructive ways"),
        ],
        next_page(),
    )

    # 10 -- real prompts: OCR text
    add_prompt_slide(
        prs, "The Actual Prompts", "Prompt 1 of 3 — Plain Text / Title / Caption / Formula",
        "vlm_client.py :: ocr_text()",
        "\"Transcribe all text in this image exactly as written. "
        "Return plain text only, no commentary, no markdown formatting.\"",
        "Used for every text-bearing region class -- title, plain text, captions, and formulas all go through this one call.",
        next_page(),
    )

    # 11 -- real prompts: table to html
    add_prompt_slide(
        prs, "The Actual Prompts", "Prompt 2 of 3 — Table → HTML",
        "vlm_client.py :: table_to_html()",
        "\"Convert this table image into a single valid HTML <table> element, "
        "preserving all rows, columns and header cells. "
        "Return only the HTML markup, no markdown code fences, no commentary.\"",
        "The output gets embedded directly into the markdown as raw HTML -- so it renders as an actual table, not a code block.",
        next_page(),
    )

    # 12 -- real prompts: figure classify
    add_prompt_slide(
        prs, "The Actual Prompts", "Prompt 3 of 3 — Figure Classify + Describe",
        "vlm_client.py :: describe_figure()",
        "\"This image is cropped from a document page. "
        "Classify it as exactly one of: chart, photo, diagram, logo, signature, stamp_seal, other. "
        "Then write a one-sentence description of what it shows. "
        "Reply in exactly this format:\nTYPE: <type>\nDESCRIPTION: <description>\"",
        "One call does double duty: the TYPE decides whether this becomes body content (chart/photo/diagram) or metadata-only (logo/signature/stamp_seal).",
        next_page(),
    )

    # 13 -- fun fact / real bug: models don't always listen
    add_fact_slide(
        prs, "Real Finding", "Small Models Don't Always Follow Instructions",
        [
            (0, "The table-to-HTML prompt explicitly says: “no markdown code fences”"),
            (0, "Qwen3.5-0.8B wrapped its answer in ```html ... ``` anyway, on a real run"),
            (0, "The pipeline now strips the fence defensively if the model adds one -- but the lesson is the point:"),
            (1, "Never treat a prompt instruction as load-bearing. Verify the output, don't just trust the request."),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 14 -- step 5 region routing
    add_bullet_slide(
        prs, "Step 5 of 7", "Region Routing — Not Everything Becomes Body Text",
        [
            (0, "title → markdown heading (##)"),
            (0, "plain text / captions / formulas → transcribed text"),
            (0, "table → raw HTML, embedded inline"),
            (0, "figure, classified as chart / photo / diagram → image + description in the body"),
            (0, "figure, classified as logo / signature / stamp_seal → EXCLUDED from the body"),
            (1, "These are page furniture and provenance markers, not document content -- they go to metadata.json instead"),
            (0, "abandon (headers/footers/watermarks) → dropped from the body, but still logged"),
        ],
        next_page(),
    )

    # 15 -- step 6 assembly + known limitation
    add_bullet_slide(
        prs, "Step 6 of 7", "Assembly — Reading Order",
        [
            (0, "Regions are sorted top-to-bottom, then left-to-right, and joined into one markdown file per document"),
            (0, "Known limitation: this is a single-column assumption"),
            (1, "On a genuinely two-column document (tested: a Census Bureau statistical brief), the naive sort interleaves the two columns and the text comes out scrambled"),
            (1, "Fixing it means detecting columns before sorting, not after -- a real extension exercise, not a bug we papered over"),
            (0, "This is exactly why the sample documents in input_pdfs/ were deliberately chosen to be single-column"),
        ],
        next_page(),
    )

    # 16 -- step 7 postprocessing + regions.json
    add_bullet_slide(
        prs, "Step 7 of 7", "Post-Processing, Validation & Metadata",
        [
            (0, "Whitespace cleanup on the final markdown"),
            (0, "Every <table> block gets a tag-balance check -- malformed ones are flagged, not silently shipped"),
            (0, "metadata.json records: native PDF metadata, region counts, artifacts detected, low-confidence regions, dropped regions"),
            (0, "regions.json records the FULL per-region layout data: true class, confidence score, page, bbox, clean text"),
            (1, "This turned out to matter beyond this module -- it's what Module 3 (Chunking) reads to make real structure-aware decisions, instead of re-guessing structure from markdown formatting"),
        ],
        next_page(),
    )

    # 17 -- device handling
    add_bullet_slide(
        prs, "Engineering Detail", "Running on Real Hardware, Not Just “a GPU”",
        [
            (0, "Every model call auto-picks a backend: CUDA → Apple MPS → CPU"),
            (0, "Requesting a device that isn't there fails immediately and clearly, rather than crashing deep inside a tensor op"),
            (0, "If an operation isn't implemented on the chosen backend at runtime (this happens on MPS for some kernels), the pipeline logs a warning and retries on CPU instead of dying mid-batch"),
            (0, "A --device flag lets you force a specific backend to compare speed and behavior directly"),
        ],
        next_page(),
    )

    # 18 -- image slide: duplicate box bug on federal register
    add_image_slide(
        prs, "Real Bug", "The Same Text, Boxed Twice",
        annotated_fedreg,
        [
            (0, "Look at “A Proclamation”: it gets BOTH a red title box (0.91 confidence) and an overlapping blue plain-text box (0.98)"),
            (0, "Result: that line is transcribed twice in the markdown output"),
            (0, "No cross-class overlap suppression (IoU-based NMS) is implemented -- this is a real, open gap, not a hidden one"),
            (1, "Good candidate for a class exercise: how would you deduplicate overlapping detections across different classes?"),
        ],
        next_page(),
        color=ACCENT_BAD,
    )

    # 19 -- fun fact: chart hallucination
    add_fact_slide(
        prs, "Real Finding", "Small Models Can Be Confidently Wrong",
        [
            (0, "Asked to describe a bar chart titled “Q3 Regional Output” with values 40–80:"),
            (1, "The model's description said “Q5” and “40 to 80,000 units”"),
            (0, "The general gist was right -- a bar chart, four regions. The specific numbers were confidently invented."),
            (0, "This is a well-known small-VLM failure mode on data-dense images: good at gist, unreliable at precise numbers"),
            (1, "Worth remembering before trusting any VLM to read a chart's axis values unsupervised"),
        ],
        next_page(),
    )

    # 20 -- wrap-up
    add_bullet_slide(
        prs, "Wrap-Up", "What Module 1 Actually Produces",
        [
            (0, "Per document: pages/ (rendered), annotated/ (boxes), images/ (cropped figures), <doc>.md, metadata.json, regions.json"),
            (0, "Every step is inspectable -- nothing is a black box you just have to trust"),
            (0, "A small, local, open-weight VLM (0.8B params) handled real government PDFs end to end"),
            (0, "Its real mistakes -- duplicate boxes, chart hallucination, a misclassified logo -- are as instructive as its successes"),
            (0, "Next: Module 2 turns this output into a validation report; Module 3 turns it into retrieval-ready chunks"),
        ],
        next_page(),
    )

    # 21 -- closing
    add_closing_slide(
        prs,
        title="Questions?",
        subtitle="1. OCR/  --  README.md has the full pipeline detail, every prompt, and every finding in this deck.",
        footer="Module 1 of 3  ·  Data Processing → Analysis → Chunking",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
